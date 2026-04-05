#!/usr/bin/env python3
"""Generate the Trinity-CAM (GPT) stakeholder presentation."""

from __future__ import annotations

from dataclasses import dataclass
from datetime import date, datetime
from pathlib import Path
import re

from openpyxl import load_workbook
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


HERE = Path(__file__).parent
PROJECT_NAME = "Trinity-CAM (GPT)"
OUTPUT_FOLDER_NAME = "Trinity-CAM (GPT) v1.1"
DOCUMENT_VERSION = "Version 1.1 - Change Request 1"
REPORT_DATE = date(2026, 4, 5)
TEMPLATE = HERE / "Reference" / "Bosch presentation template.pptx"
LOGO = HERE / "Bosch.png"
OUTPUT_DIR = HERE / "active-projects" / OUTPUT_FOLDER_NAME
OUTPUT = OUTPUT_DIR / f"{PROJECT_NAME}_Stakeholder_Presentation.pptx"

SCHEDULE_PATH = OUTPUT_DIR / f"{PROJECT_NAME}_Project_Schedule.xlsx"
RISK_PATH = OUTPUT_DIR / f"{PROJECT_NAME}_Risk_Register.xlsx"
COST_PATH = OUTPUT_DIR / f"{PROJECT_NAME}_Cost_Plan.xlsx"

BLUE = RGBColor(0x00, 0x62, 0xA3)
BLUE_DARK = RGBColor(0x00, 0x3B, 0x5C)
BLUE_LIGHT = RGBColor(0xD9, 0xEE, 0xF7)
RED = RGBColor(0xE2, 0x00, 0x15)
GREEN = RGBColor(0x27, 0x7A, 0x3E)
AMBER = RGBColor(0xC7, 0x76, 0x00)
GREY = RGBColor(0x66, 0x66, 0x66)
GREY_LIGHT = RGBColor(0xF2, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x11, 0x11, 0x11)


IMPACT_SCORES = {
    "Very Low": 1,
    "Low": 2,
    "Moderate": 3,
    "High": 4,
    "Very High": 5,
}

PROBABILITY_SCORES = {
    "10%": 1,
    "30%": 2,
    "50%": 3,
    "70%": 4,
    "90%": 5,
}


@dataclass
class Milestone:
    name: str
    when: date


@dataclass
class NearTermTask:
    name: str
    start: date
    finish: date
    note: str


@dataclass
class RiskItem:
    risk_id: str
    category: str
    event: str
    owner: str
    strategy: str
    status: str
    impact: str
    probability: str
    score: int


def as_date(value: object) -> date:
    if isinstance(value, datetime):
        return value.date()
    if isinstance(value, date):
        return value
    raise TypeError(f"Unsupported date value: {value!r}")


def money_compact(amount: int) -> str:
    return f"EUR {amount / 1_000_000:.2f}M"


def format_date(value: date) -> str:
    return value.strftime("%d %b %Y")


def clear_template_slides(prs: Presentation) -> None:
    slide_ids = prs.slides._sldIdLst
    for slide_id in list(slide_ids):
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        slide_ids.remove(slide_id)


def set_shape_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = str(text)
    paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, lines, font_size=18, color=BLACK):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = str(line)
        paragraph.level = 0
        paragraph.bullet = True
        paragraph.space_after = Pt(9)
        run = paragraph.runs[0]
        run.font.name = "Aptos"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return box


def add_branding(slide, title: str, subtitle: str) -> None:
    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.6))
    set_shape_fill(header, BLUE)
    header.line.fill.background()

    footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(7.1), Inches(13.333), Inches(0.4))
    set_shape_fill(footer, BLUE_DARK)
    footer.line.fill.background()

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(1.18), Inches(0.12), Inches(5.45))
    set_shape_fill(accent, RED)
    accent.line.fill.background()

    add_textbox(slide, Inches(0.75), Inches(0.72), Inches(9.2), Inches(0.4), title, 24, BLUE_DARK, bold=True)
    add_textbox(slide, Inches(0.77), Inches(1.03), Inches(10.0), Inches(0.28), subtitle, 11, GREY)
    add_textbox(slide, Inches(0.55), Inches(7.13), Inches(4.6), Inches(0.2), "Robert Bosch GmbH | Trinity-CAM carve-out", 9, WHITE)
    add_textbox(slide, Inches(10.25), Inches(7.13), Inches(2.5), Inches(0.2), REPORT_DATE.strftime("Report date %d %b %Y"), 9, WHITE, align=PP_ALIGN.RIGHT)

    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(11.35), Inches(0.08), height=Inches(0.38))


def add_card(slide, left, top, width, height, headline, value, fill, value_color=BLUE_DARK):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    set_shape_fill(card, fill)
    card.line.color.rgb = WHITE
    add_textbox(slide, left + Inches(0.12), top + Inches(0.10), width - Inches(0.24), Inches(0.25), headline, 10, GREY)
    add_textbox(slide, left + Inches(0.12), top + Inches(0.32), width - Inches(0.24), Inches(0.42), value, 20, value_color, bold=True)


def load_schedule():
    sheet = load_workbook(SCHEDULE_PATH, data_only=True)["Schedule"]
    milestones = []
    near_term = []
    for row in sheet.iter_rows(min_row=2, values_only=True):
        name = (row[2] or "").strip()
        start = row[4]
        finish = row[5]
        notes = row[8] or ""
        is_milestone = row[9] == "Yes"
        if start and finish and is_milestone:
            milestones.append(Milestone(name=name, when=as_date(start)))
        if start and finish and row[1] == 3 and not is_milestone and as_date(start) <= date(2026, 9, 30):
            near_term.append(
                NearTermTask(name=name, start=as_date(start), finish=as_date(finish), note=notes)
            )
    return milestones, near_term[:6]


def load_risks():
    sheet = load_workbook(RISK_PATH, data_only=True)["Risk Register"]
    risks = []
    for row_idx in range(5, sheet.max_row + 1):
        risk_id = sheet.cell(row_idx, 2).value
        if not isinstance(risk_id, str) or not re.fullmatch(r"R\d{3}", risk_id):
            continue
        impact = sheet.cell(row_idx, 28).value or "Moderate"
        probability = sheet.cell(row_idx, 30).value or "30%"
        score = IMPACT_SCORES.get(impact, 3) * PROBABILITY_SCORES.get(probability, 2)
        risks.append(
            RiskItem(
                risk_id=risk_id,
                category=sheet.cell(row_idx, 4).value or "",
                event=sheet.cell(row_idx, 6).value or "",
                owner=sheet.cell(row_idx, 9).value or "Unassigned",
                strategy=sheet.cell(row_idx, 22).value or "Mitigate",
                status=(sheet.cell(row_idx, 26).value or "not started").title(),
                impact=impact,
                probability=probability,
                score=score,
            )
        )
    risks.sort(key=lambda item: (item.score, item.risk_id), reverse=True)
    return risks


def load_costs():
    sheet = load_workbook(COST_PATH, data_only=True)["Cost Plan"]
    categories = []
    for row_idx in range(81, 91):
        label = sheet.cell(row_idx, 1).value
        amount = sheet.cell(row_idx, 6).value
        if label and amount is not None:
            categories.append((str(label), int(amount)))

    phases = []
    for row_idx in range(93, 99):
        label = sheet.cell(row_idx, 1).value
        amount = sheet.cell(row_idx, 6).value
        if label and amount is not None:
            phases.append((str(label), int(amount)))

    capex_lines = []
    for row_idx in range(101, 108):
        label = sheet.cell(row_idx, 1).value
        estimate = sheet.cell(row_idx, 5).value
        if label and estimate:
            capex_lines.append((str(label), str(estimate)))

    labour_total = sum(amount for _, amount in phases)
    budget_note = str(sheet.cell(114, 1).value or "")
    return categories, phases, capex_lines, labour_total, budget_note


def add_cover_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[19])
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.35))
    set_shape_fill(band, BLUE)
    band.line.fill.background()

    lower = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(6.65), Inches(13.333), Inches(0.85))
    set_shape_fill(lower, BLUE_DARK)
    lower.line.fill.background()

    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.75), Inches(1.85), Inches(0.16), Inches(3.65))
    set_shape_fill(stripe, RED)
    stripe.line.fill.background()

    add_textbox(slide, Inches(1.12), Inches(1.75), Inches(8.8), Inches(0.8), PROJECT_NAME, 28, BLUE_DARK, bold=True)
    add_textbox(slide, Inches(1.14), Inches(2.45), Inches(8.7), Inches(0.45), "Management Stakeholder Presentation", 20, BLUE)
    add_textbox(slide, Inches(1.14), Inches(3.05), Inches(9.3), Inches(1.05), "Integration model carve-out from JCI to Bosch via an Infosys-operated merger zone", 24, BLACK, bold=True)
    add_textbox(slide, Inches(1.14), Inches(4.45), Inches(7.8), Inches(0.3), f"Report date: {REPORT_DATE.strftime('%d %B %Y')} | {DOCUMENT_VERSION}", 14, GREY)
    add_textbox(slide, Inches(1.14), Inches(4.8), Inches(8.8), Inches(0.5), "48 sites | 12,000 IT users | 1,800+ applications | GoLive 01 Jan 2028", 16, GREY)
    add_textbox(slide, Inches(0.75), Inches(6.88), Inches(6.5), Inches(0.2), "Seller: Johnson Controls International | Buyer: Robert Bosch GmbH", 10, WHITE)
    add_textbox(slide, Inches(9.05), Inches(6.88), Inches(3.55), Inches(0.2), "Prepared from generated schedule, risk, cost, and status outputs", 9, WHITE, align=PP_ALIGN.RIGHT)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(10.8), Inches(0.25), height=Inches(0.55))


def add_summary_slide(prs: Presentation, milestones, labour_total, risks):
    slide = prs.slides.add_slide(prs.slide_layouts[19])
    add_branding(slide, "Executive Summary", "April 2026 change-request view with approved seller-side buffer through 31 Jul 2027")

    add_bullets(
        slide,
        Inches(0.95),
        Inches(1.55),
        Inches(7.0),
        Inches(3.0),
        [
            "The carve-out transfers the JCI air conditioning business into Bosch using a temporary merger zone operated by Infosys.",
            "All 12,000 users and 1,800+ applications start on the JCI side; the programme therefore depends on a clean JCI to merger-zone to Bosch transition path.",
            "JCI has approved a TSA extension through 31 Jul 2027 so users can remain in the legacy environment while Infosys continues merger-zone build-up.",
            f"The derived external labour baseline is {money_compact(labour_total)}; CAPEX and contingency decisions remain gated to QG1.",
            f"The current register holds {len(risks)} threats with the highest exposure concentrated in SAP critical path, merger-zone readiness, and pre-GoLive quality gates.",
        ],
    )

    add_card(slide, Inches(8.25), Inches(1.62), Inches(2.05), Inches(1.0), "Programme start", format_date(milestones[0].when), BLUE_LIGHT)
    add_card(slide, Inches(10.45), Inches(1.62), Inches(2.05), Inches(1.0), "GoLive", format_date(milestones[4].when), BLUE_LIGHT)
    add_card(slide, Inches(8.25), Inches(2.88), Inches(2.05), Inches(1.0), "Completion", format_date(milestones[5].when), BLUE_LIGHT)
    add_card(slide, Inches(10.45), Inches(2.88), Inches(2.05), Inches(1.0), "Top risk score", str(risks[0].score), RGBColor(0xFD, 0xE7, 0xE9), value_color=RED)

    note = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.25), Inches(4.25), Inches(4.25), Inches(1.85))
    set_shape_fill(note, GREY_LIGHT)
    note.line.color.rgb = WHITE
    add_textbox(slide, Inches(8.48), Inches(4.45), Inches(3.8), Inches(0.25), "Management message", 11, GREY, bold=True)
    add_textbox(slide, Inches(8.48), Inches(4.78), Inches(3.7), Inches(1.1), "This is a pre-kickoff control baseline rather than an in-flight status deck. The immediate focus is governance stand-up, Infosys onboarding, and QG1 concept maturity.", 15, BLUE_DARK)


def add_scope_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[19])
    add_branding(slide, "Programme Scope & Operating Model", "Integration carve-out design for Trinity-CAM")

    model_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(1.6), Inches(12.0), Inches(1.1))
    set_shape_fill(model_box, BLUE_LIGHT)
    model_box.line.color.rgb = WHITE
    add_textbox(slide, Inches(1.18), Inches(1.82), Inches(11.5), Inches(0.28), "JCI source IT -> Infosys merger zone -> Bosch target IT", 21, BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.18), Inches(2.15), Inches(11.5), Inches(0.22), "The merger zone is a temporary operational bridge, not the target-state platform.", 11, GREY, align=PP_ALIGN.CENTER)

    left = [
        "Seller: Johnson Controls International air conditioning business.",
        "Buyer: Robert Bosch GmbH as sponsor customer and target operating owner.",
        "PMO and methodology lead: KPMG.",
        "Delivery partner: Infosys for merger-zone setup, run, IT services, and migrations.",
    ]
    right = [
        "Global scope covers 48 sites and 12,000 IT users.",
        "Application estate exceeds 1,800 systems including a major SAP landscape.",
        "TSA now supports continuity through 31 Jul 2027 and removes near-term pressure from Bosch without changing the overall programme timeline.",
        "GoLive marks Day 1 operation in the merger zone; Bosch steady-state handover completes by QG5.",
    ]
    add_bullets(slide, Inches(1.0), Inches(3.1), Inches(5.45), Inches(2.75), left, 17)
    add_bullets(slide, Inches(6.7), Inches(3.1), Inches(5.45), Inches(2.75), right, 17)

    add_card(slide, Inches(1.0), Inches(6.0), Inches(2.25), Inches(0.78), "Sites", "48", BLUE_LIGHT)
    add_card(slide, Inches(3.45), Inches(6.0), Inches(2.25), Inches(0.78), "Users", "12,000", BLUE_LIGHT)
    add_card(slide, Inches(5.9), Inches(6.0), Inches(2.25), Inches(0.78), "Applications", "1,800+", BLUE_LIGHT)
    add_card(slide, Inches(8.35), Inches(6.0), Inches(4.0), Inches(0.78), "Delivery model", "Integration via temporary merger zone", BLUE_LIGHT)


def add_timeline_slide(prs: Presentation, milestones, near_term):
    slide = prs.slides.add_slide(prs.slide_layouts[19])
    add_branding(slide, "Integrated Timeline & Stage Gates", "Dates are read directly from the generated programme schedule")

    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.1), Inches(2.3), Inches(11.0), Inches(0.08))
    set_shape_fill(line, BLUE)
    line.line.fill.background()

    milestone_positions = [1.15, 3.2, 5.35, 7.55, 9.5, 11.55]
    for xpos, milestone in zip(milestone_positions, milestones):
        marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(xpos), Inches(2.05), Inches(0.28), Inches(0.28))
        set_shape_fill(marker, RED if "GoLive" in milestone.name else BLUE_DARK)
        marker.line.fill.background()
        add_textbox(slide, Inches(xpos - 0.15), Inches(2.5), Inches(1.5), Inches(0.42), milestone.name.replace("    ", " "), 10, BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(xpos - 0.08), Inches(3.0), Inches(1.35), Inches(0.24), format_date(milestone.when), 10, GREY, align=PP_ALIGN.CENTER)

    phase_note = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(3.75), Inches(11.95), Inches(0.72))
    set_shape_fill(phase_note, GREY_LIGHT)
    phase_note.line.color.rgb = WHITE
    add_textbox(slide, Inches(1.15), Inches(3.95), Inches(11.5), Inches(0.28), "Phase sequence: Mobilize -> Discover/Design -> Build -> Test/Migrate -> Final Readiness -> Hypercare/Handover", 15, BLUE_DARK, bold=True, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.0), Inches(4.8), Inches(3.9), Inches(0.28), "Next 90 days to QG1", 15, BLUE_DARK, bold=True)
    task_lines = [
        f"{format_date(task.start)} to {format_date(task.finish)}: {task.name}" for task in near_term[:4]
    ]
    add_bullets(slide, Inches(1.0), Inches(5.1), Inches(5.6), Inches(1.55), task_lines, 13)

    note = [
        "QG4 remains separated from GoLive by a controlled final-readiness buffer.",
        "All post-GoLive work is stabilization and Bosch handover only.",
        "The schedule has no summary-task predecessor defects or gate/date collisions.",
    ]
    add_bullets(slide, Inches(7.0), Inches(4.95), Inches(5.2), Inches(1.55), note, 13)


def add_budget_slide(prs: Presentation, categories, phases, capex_lines, labour_total, budget_note):
    slide = prs.slides.add_slide(prs.slide_layouts[19])
    add_branding(slide, "Budget & Cost Structure", "Derived from the generated labour cost plan and aligned to the risk register")

    add_card(slide, Inches(0.95), Inches(1.55), Inches(2.35), Inches(0.95), "External labour baseline", money_compact(labour_total), BLUE_LIGHT)
    add_card(slide, Inches(3.5), Inches(1.55), Inches(2.35), Inches(0.95), "Largest phase", money_compact(max(amount for _, amount in phases)), BLUE_LIGHT)
    add_card(slide, Inches(6.05), Inches(1.55), Inches(2.35), Inches(0.95), "Budget release gate", "QG1", BLUE_LIGHT)
    add_card(slide, Inches(8.6), Inches(1.55), Inches(3.35), Inches(0.95), "Budget note", budget_note.replace("6. ", ""), RGBColor(0xFD, 0xF5, 0xD8), value_color=AMBER)

    add_textbox(slide, Inches(1.0), Inches(2.85), Inches(4.3), Inches(0.25), "Cost by delivery team", 14, BLUE_DARK, bold=True)
    top_categories = sorted(categories, key=lambda item: item[1], reverse=True)[:5]
    max_amount = top_categories[0][1]
    bar_top = 3.2
    for index, (label, amount) in enumerate(top_categories):
        ypos = Inches(bar_top + index * 0.58)
        add_textbox(slide, Inches(1.0), ypos, Inches(2.75), Inches(0.18), label, 10, BLACK)
        width = 2.65 * amount / max_amount
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(3.8), ypos + Inches(0.03), Inches(width), Inches(0.19))
        set_shape_fill(bar, BLUE if index < 2 else BLUE_DARK)
        bar.line.fill.background()
        add_textbox(slide, Inches(6.6), ypos - Inches(0.01), Inches(1.15), Inches(0.18), money_compact(amount), 10, GREY, align=PP_ALIGN.RIGHT)

    add_textbox(slide, Inches(7.2), Inches(2.85), Inches(4.3), Inches(0.25), "Phase concentration", 14, BLUE_DARK, bold=True)
    phase_lines = [
        f"{label.split(':', 1)[0]}: {money_compact(amount)}" for label, amount in phases[1:5]
    ]
    add_bullets(slide, Inches(7.2), Inches(3.15), Inches(4.8), Inches(1.75), phase_lines, 14)

    capex_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.15), Inches(5.15), Inches(5.0), Inches(1.2))
    set_shape_fill(capex_box, GREY_LIGHT)
    capex_box.line.color.rgb = WHITE
    add_textbox(slide, Inches(7.35), Inches(5.33), Inches(4.6), Inches(0.22), "CAPEX / contingency carried outside labour baseline", 11, GREY, bold=True)
    capex_preview = [f"{label}: {estimate}" for label, estimate in capex_lines[:3]]
    add_bullets(slide, Inches(7.35), Inches(5.58), Inches(4.45), Inches(0.6), capex_preview, 10)


def add_risk_slide(prs: Presentation, risks):
    slide = prs.slides.add_slide(prs.slide_layouts[19])
    add_branding(slide, "Top Risks & Management Actions", "Prioritization derived from actual probability and impact fields in the GPT risk register")

    high = sum(1 for risk in risks if risk.score >= 15)
    medium = sum(1 for risk in risks if 10 <= risk.score < 15)
    add_card(slide, Inches(0.95), Inches(1.55), Inches(2.1), Inches(0.88), "Open risks", str(len(risks)), BLUE_LIGHT)
    add_card(slide, Inches(3.2), Inches(1.55), Inches(2.1), Inches(0.88), "High exposure", str(high), RGBColor(0xFD, 0xE7, 0xE9), value_color=RED)
    add_card(slide, Inches(5.45), Inches(1.55), Inches(2.1), Inches(0.88), "Medium exposure", str(medium), RGBColor(0xFD, 0xF5, 0xD8), value_color=AMBER)
    add_card(slide, Inches(7.7), Inches(1.55), Inches(4.25), Inches(0.88), "Control focus", "SAP, merger-zone readiness, wave execution, data quality", BLUE_LIGHT)

    top = risks[:5]
    for index, risk in enumerate(top):
        ypos = 2.8 + index * 0.78
        row = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(ypos), Inches(11.95), Inches(0.58))
        fill = RGBColor(0xFD, 0xE7, 0xE9) if risk.score >= 15 else GREY_LIGHT
        set_shape_fill(row, fill)
        row.line.color.rgb = WHITE
        add_textbox(slide, Inches(1.15), Inches(ypos + 0.09), Inches(0.65), Inches(0.18), risk.risk_id, 11, BLUE_DARK, bold=True)
        add_textbox(slide, Inches(1.9), Inches(ypos + 0.08), Inches(5.55), Inches(0.2), risk.event[:120].rstrip() + ("..." if len(risk.event) > 120 else ""), 10, BLACK)
        add_textbox(slide, Inches(7.6), Inches(ypos + 0.08), Inches(1.0), Inches(0.2), f"Score {risk.score}", 10, RED if risk.score >= 15 else BLUE_DARK, bold=True)
        add_textbox(slide, Inches(8.7), Inches(ypos + 0.08), Inches(1.55), Inches(0.2), risk.owner, 10, GREY)
        add_textbox(slide, Inches(10.45), Inches(ypos + 0.08), Inches(1.25), Inches(0.2), risk.strategy, 10, BLUE)

    add_textbox(slide, Inches(1.0), Inches(6.75), Inches(9.8), Inches(0.22), "Management direction: hold a weekly integrated risk review from QG1 onward, with SAP and merger-zone controls tracked as explicit steering actions rather than buried in workstream logs.", 11, BLUE_DARK)


def add_decisions_slide(prs: Presentation, milestones):
    slide = prs.slides.add_slide(prs.slide_layouts[19])
    add_branding(slide, "Management Decisions Required", "Actions to keep the programme viable ahead of QG1 and QG4")

    decisions = [
        "Approve the QG1 design and funding gate for merger-zone architecture, security controls, and programme governance on 01 Oct 2026.",
        "Confirm Infosys as the accountable delivery partner for merger-zone build, operation, and migration throughput targets.",
        "Mandate JCI data, application, and TSA tower owners to deliver named counterparts and decision SLAs before mobilisation completes.",
        "Confirm Bosch target-state integration resources early enough to absorb the Day 1 merger-zone landing and the QG5 handover path.",
        f"Maintain {format_date(milestones[4].when)} as the hard Day 1 date and treat the {format_date(milestones[3].when)} QG4 gate as non-negotiable.",
    ]
    add_bullets(slide, Inches(1.0), Inches(1.7), Inches(7.35), Inches(4.5), decisions, 18)

    right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.8), Inches(3.15), Inches(4.35))
    set_shape_fill(right, GREY_LIGHT)
    right.line.color.rgb = WHITE
    add_textbox(slide, Inches(9.02), Inches(2.05), Inches(2.7), Inches(0.28), "Key dates", 13, GREY, bold=True)
    add_textbox(slide, Inches(9.02), Inches(2.42), Inches(2.7), Inches(1.9), "QG1: 01 Oct 2026\nApproved JCI TSA end: 31 Jul 2027\nQG2/QG3: 31 Jul 2027\nQG4: 10 Dec 2027\nGoLive: 01 Jan 2028\nQG5: 01 Apr 2028", 16, BLUE_DARK)
    add_textbox(slide, Inches(9.02), Inches(4.65), Inches(2.6), Inches(1.0), "Programme principle:\nNo hidden scope transfer from JCI to Bosch without governance, budget, and exit evidence.", 12, BLUE)


def main() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    milestones, near_term = load_schedule()
    risks = load_risks()
    categories, phases, capex_lines, labour_total, budget_note = load_costs()

    prs = Presentation(TEMPLATE)
    clear_template_slides(prs)

    add_cover_slide(prs)
    add_summary_slide(prs, milestones, labour_total, risks)
    add_scope_slide(prs)
    add_timeline_slide(prs, milestones, near_term)
    add_budget_slide(prs, categories, phases, capex_lines, labour_total, budget_note)
    add_risk_slide(prs, risks)
    add_decisions_slide(prs, milestones)

    prs.save(str(OUTPUT))
    print(f"[Trinity-CAM (GPT)] Stakeholder Presentation: {OUTPUT}")


if __name__ == "__main__":
    main()