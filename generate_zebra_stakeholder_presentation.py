#!/usr/bin/env python3
"""
Generate Zebra_Stakeholder_Presentation.pptx

Final executive presentation for Project Zebra stakeholders.
Uses Bosch presentation template as base.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

HERE = Path(__file__).parent
TEMPLATE_PATH = HERE / "Bosch presentation template.pptx"
OUTPUT_PATH = HERE / "active-projects" / "Zebra" / "Zebra_Stakeholder_Presentation.pptx"

# Ensure output directory exists
OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)

# Bosch colours
DARK_BLUE = RGBColor(0, 59, 110)        # #003B6E
MID_BLUE = RGBColor(0, 102, 204)        # #0066CC
LIGHT_BLUE = RGBColor(228, 237, 249)    # #E4EDFF
GREY = RGBColor(95, 95, 95)             # #5F5F5F

def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Format
    title_shape.text_frame.paragraphs[0].font.size = Pt(54)
    title_shape.text_frame.paragraphs[0].font.bold = True
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
    
    return slide

def add_content_slide(prs, title, subtitle_text=None):
    """Add a content slide with title and optional subtitle."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    if subtitle_text:
        subtitle_shape = slide.placeholders[1]
        tf = subtitle_shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = subtitle_text
        p.font.size = Pt(20)
        p.font.color.rgb = GREY
    
    return slide

def add_bullet_slide(prs, title, bullets):
    """Add a slide with bullet points."""
    slide = add_content_slide(prs, title)
    content_placeholder = slide.placeholders[1]
    tf = content_placeholder.text_frame
    tf.clear()
    
    for bullet_text in bullets:
        p = tf.add_paragraph()
        p.text = bullet_text
        p.level = 0
        p.font.size = Pt(16)
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    return slide

def main():
    print("[Zebra] Generating stakeholder presentation...")
    
    # Load template or create new presentation
    try:
        prs = Presentation(TEMPLATE_PATH)
        print(f"  Template loaded from {TEMPLATE_PATH}")
    except FileNotFoundError:
        print(f"  Template not found, creating presentation from scratch")
        prs = Presentation()
    
    # Clear default slides if any
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    # Slide 1: Cover
    add_title_slide(prs, "PROJECT ZEBRA", "Packaging Business Carve-Out | Stand Alone Model | 20-Month Journey")
    
    # Slide 2: Project Overview
    add_bullet_slide(prs, "Project Overview", [
        "Objective: Complete separation of Packaging business from Bosch global IT infrastructure",
        "Scope: 37 worldwide sites, 3,500+ IT users, 208 applications, SAP ERP",
        "Timeline: 1 April 2026 – 31 October 2027 (20 months)",
        "GoLive: 1 June 2027",
        "Model: Stand Alone with Transition Services Agreement (TSA) support post-GoLive",
        "Budget: EUR 2.34M (labour + CAPEX)",
    ])
    
    # Slide 3: Carve-Out Model & Key Parties
    add_bullet_slide(prs, "Carve-Out Model & Governance", [
        "Stand Alone Model: Buyer assumes 100% operational control at GoLive; no shared infrastructure post-GoLive",
        "Seller: Robert Bosch GmbH (carve-out entity)",
        "Buyer: [Confidential — legal hold]",
        "Programme Manager: Gill Amandeep Singh (BD/MIL-PSM1)",
        "PMO Partner: KPMG",
        "TSA: Seller IT support for 4-6 months post-GoLive to ease transition",
    ])
    
    # Slide 4: Timeline & Key Milestones
    add_bullet_slide(prs, "Timeline & Quality Gates", [
        "Phase 0 (Initialization): April 1–17, 2026 → QG0 (17 April)",
        "Phase 1 (Concept): April 18 – May 18, 2026 → QG1 (18 May)",
        "Phase 2 (Design & Architecture): May 19 – July 27, 2026 → QG2/3 (27 July)",
        "Phase 3 (Build & Test): July 28, 2026 – May 27, 2027",
        "Phase 4 (GoLive & Closure): June 1 – October 31, 2027 → QG4 (Pre-GoLive) & QG5 (Closure)",
        "Critical Path: SAP separation, 37-site multi-geography cutover, 208-app portfolio transition",
    ])
    
    # Slide 5: Scope & Complexity
    add_bullet_slide(prs, "Scope & Complexity Drivers", [
        "37 Worldwide Sites: Global multi-geography footprint (EMEA, Americas, APAC) requiring coordinated cutover",
        "3,500+ IT Users: Large-scale training, change management, and user acceptance testing",
        "208 Applications: Diverse portfolio (SaaS, on-prem, legacy); ISV licensing transitions; vendor negotiations",
        "SAP ERP: Mission-critical system requiring careful data separation and business continuity controls",
        "Data Residency & Compliance: GDPR, local data protection, export controls across 37 sites",
        "Stand Alone Requirement: No residual seller systems post-GoLive; 100% buyer IT independence by Day 1",
    ])
    
    # Slide 6: Budget Breakdown
    add_bullet_slide(prs, "Budget Summary: EUR 2.34M", [
        "Labour (KPMG + Seller IT): EUR 2.12M",
        "  • KPMG PMO & Governance: EUR 285K",
        "  • KPMG SAP & ERP Build: EUR 630K",
        "  • KPMG Testing & QA: EUR 495K",
        "  • KPMG Change Mgmt: EUR 280K",
        "  • KPMG Data & Integration: EUR 330K",
        "  • RoboGmbH IT Support: EUR 228K",
        "CAPEX: EUR 226K (SAP audits, app tools, compliance, contingency)",
        "Status: Baseline approved at QG0",
    ])
    
    # Slide 7: Top Risks
    add_bullet_slide(prs, "Top 5 Risks (P×I ≥ 12)", [
        "🔴 Multi-Site Cutover Coordination (P5×I5=25): Parallel rollout across 37 sites; single site failure impacts global ops",
        "🔴 SAP Data Separation (P4×I4=16): Packaging/other BU data leakage post-GoLive; regulatory exposure",
        "🔴 Buyer Infrastructure Readiness (P4×I4=16): Buyer IT capacity unknown due to confidentiality; readiness gap by GoLive",
        "🟡 208-App Portfolio Transitions (P4×I4=16): ISV licensing disputes; vendor lock-in; 15–20 apps at risk",
        "🟡 Regulatory Compliance (P3×I5=15): GDPR + data residency rules; global scope amplifies complexity",
    ])
    
    # Slide 8: Workstream Coverage
    add_bullet_slide(prs, "IT Workstream Coverage (9 Workstreams)", [
        "1. SAP ERP Separation | 2. Application Portfolio (208 Apps) | 3. Data Migration",
        "4. Infrastructure & Network | 5. Security & IAM | 6. Testing & QA",
        "7. Change Management & Training | 8. TSA & Handover | 9. Programme Management",
        "Key Dependencies: Phase 2 designs must be finalized by late July to allow adequate Phase 3 build/test time",
        "Resource Confidence: 85% assigned; remaining 15% of specialized SAP/infrastructure roles ramping May–June",
    ])
    
    # Slide 9: Critical Success Factors
    add_bullet_slide(prs, "Critical Success Factors", [
        "✓ Buyer IT team fully engaged and infrastructure roadmap locked by end of Phase 2 (July 27)",
        "✓ SAP separation logic validated in 3 dry-run migrations before cutover (no data loss, no leakage)",
        "✓ All 37 sites pass pre-cutover readiness checks; local teams certified on buyer systems by June 1",
        "✓ 208-application portfolio ISV licensing resolved; all apps deployed in buyer infrastructure by May 27",
        "✓ Regulatory compliance (GDPR, data residency, export controls) validated pre-GoLive",
        "✓ TSA service levels and exit criteria formally agreed in writing by May 18, 2026",
    ])
    
    # Slide 10: Next Steps (90-Day Outlook)
    add_bullet_slide(prs, "Next Steps & 90-Day Outlook", [
        "Week 1-2 (Apr 4–18): Complete Phase 0 tasks; conduct QG0 gate review",
        "Week 3-6 (Apr 19 – May 16): Commence Phase 1; AS-IS documentation; requirements capture; vendor outreach",
        "Week 7-13 (May 17 – Jun 27): Complete Phase 2 design; finalize TSA terms; prepare QG2/3 gate",
        "Critical Dependency: Buyer legal disclosure must occur by May 1 to avoid Phase 1 schedule impact",
        "Escalation Plan: If buyer disclosure delays beyond May 1, parallel design activities will be triggered",
    ])
    
    # Slide 11: Call to Action
    add_bullet_slide(prs, "Our Success Metrics", [
        "📅 Deliver on 20-month timeline: 1 April 2026 – 31 October 2027",
        "💰 Maintain budget discipline: EUR 2.34M baseline + contingency for risk mitigation",
        "✅ QG4 (Pre-GoLive): Zero critical P1 blockers; all readiness criteria met",
        "🚀 GoLive Day 1 (1 June 2027): 3,500+ users and 208 apps operational on buyer systems",
        "📊 Post-GoLive: TSA support through closure; hypercare completed by September 2027",
        "🤝 Governance: Weekly steering committee reviews; transparent risk reporting; escalation clarity",
    ])
    
    # Save presentation
    prs.save(OUTPUT_PATH)
    print(f"✓ Stakeholder presentation saved: {OUTPUT_PATH}")
    print(f"  {len(prs.slides)} slides generated")

if __name__ == "__main__":
    main()
