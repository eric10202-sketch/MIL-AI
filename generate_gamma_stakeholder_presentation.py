#!/usr/bin/env python3
"""Generate the Gamma stakeholder presentation as a PPTX deck."""

from __future__ import annotations

from datetime import datetime
from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


HERE = Path(__file__).parent
PROJECT_NAME = "Gamma"
OUTPUT_FOLDER_NAME = "Gamma v1.0"
SELLER = "Robert Bosch China"
BUYER = "Alibaba"
MODEL = "Combination"
BUSINESS = "Bosch Cloud business"
REPORT_DATE = datetime(2026, 4, 5)

TEMPLATE_PATH = HERE / "Reference" / "Bosch presentation template.pptx"
SCHEDULE_PATH = HERE / "active-projects" / OUTPUT_FOLDER_NAME / f"{PROJECT_NAME}_Project_Schedule.xlsx"
RISK_PATH = HERE / "active-projects" / OUTPUT_FOLDER_NAME / f"{PROJECT_NAME}_Risk_Register.xlsx"
COST_PATH = HERE / "active-projects" / OUTPUT_FOLDER_NAME / f"{PROJECT_NAME}_Cost_Plan.xlsx"
OUTPUT_PATH = HERE / "active-projects" / OUTPUT_FOLDER_NAME / f"{PROJECT_NAME}_Stakeholder_Presentation.pptx"
LOGO_PATH = HERE / "Bosch.png"

BLUE = RGBColor(0, 59, 110)
MID_BLUE = RGBColor(0, 102, 204)
LIGHT = RGBColor(228, 237, 249)
DARK = RGBColor(26, 26, 26)
GREY = RGBColor(95, 109, 131)
RED = RGBColor(204, 0, 0)
AMBER = RGBColor(232, 160, 0)
GREEN = RGBColor(0, 122, 51)

IMPACT_VALUES = {"Very Low": 1, "Low": 2, "Moderate": 3, "High": 4, "Very High": 5}
PROBABILITY_VALUES = {"10%": 1, "30%": 2, "50%": 3, "70%": 4, "90%": 5}


def fmt_date(value) -> str:
    if hasattr(value, "strftime"):
        return value.strftime("%d %b %Y")
    return str(value)


def load_schedule_summary():
    wb = load_workbook(SCHEDULE_PATH, data_only=False)
    ws = wb["Schedule"]
    milestones = []
    for row in range(2, ws.max_row + 1):
        if ws.cell(row, 2).value is None:
            continue
        name = str(ws.cell(row, 3).value or "").strip()
        if str(ws.cell(row, 10).value) == "Yes" and ("QG" in name or "GoLive" in name or "Closure" in name):
            milestones.append((name, fmt_date(ws.cell(row, 5).value)))
    return milestones


def load_risk_summary():
    wb = load_workbook(RISK_PATH, data_only=False)
    ws = wb["Risk Register"]
    risks = []
    for row in range(5, 140):
        risk_id = ws.cell(row, 2).value
        if risk_id is None:
            continue
        impact = ws.cell(row, 12).value
        probability = ws.cell(row, 14).value
        score = IMPACT_VALUES.get(impact, 0) * PROBABILITY_VALUES.get(probability, 0)
        if ws.cell(row, 16).value == "threat":
            risks.append((int(risk_id), score, ws.cell(row, 4).value, ws.cell(row, 6).value))
    risks.sort(key=lambda item: item[1], reverse=True)
    return risks[:5]


def load_budget_summary():
    wb = load_workbook(COST_PATH, data_only=False)
    ws = wb["Cost Plan"]
    labour_total = 0
    capex_total = 0
    categories = []
    for row in range(1, ws.max_row + 1):
        label = ws.cell(row, 1).value
        if label == "OVERALL PROJECT TOTAL - LABOUR ONLY":
            labour_total = int(ws.cell(row, 6).value or 0)
        if label == "TOTAL CAPEX / ADDITIONAL COSTS":
            capex_total = int(ws.cell(row, 6).value or 0)
        if isinstance(label, str) and label.startswith("SUBTOTAL -"):
            categories.append((label.replace("SUBTOTAL - ", ""), int(ws.cell(row, 6).value or 0)))
    categories.sort(key=lambda item: item[1], reverse=True)
    return labour_total, capex_total, categories[:4]


def clear_slides(prs: Presentation) -> None:
    while len(prs.slides) > 0:
        rel_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rel_id)
        del prs.slides._sldIdLst[0]


def add_branding(slide, title: str, subtitle: str = ""):
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.65)).fill.solid()
    header = slide.shapes[-1]
    header.fill.fore_color.rgb = BLUE
    header.line.color.rgb = BLUE

    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(7.5), Inches(0.35))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.72), Inches(8.5), Inches(0.35))
        sp = sub_box.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(12)
        sp.font.color.rgb = GREY

    slide.shapes.add_picture(str(LOGO_PATH), Inches(11.45), Inches(0.09), height=Inches(0.42))
    footer = slide.shapes.add_textbox(Inches(0.45), Inches(7.0), Inches(12.0), Inches(0.25))
    fp = footer.text_frame.paragraphs[0]
    fp.text = f"{PROJECT_NAME} | {REPORT_DATE.strftime('%d %b %Y')} | Confidential"
    fp.font.size = Pt(10)
    fp.font.color.rgb = GREY


def add_bullets(slide, bullets, left=0.6, top=1.3, width=12.0, height=5.2, font_size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for index, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)


def add_table_slide(prs, title, subtitle, headers, rows, col_widths):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_branding(slide, title, subtitle)
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(0.5), Inches(1.4), Inches(12.2), Inches(4.9)).table
    for col, width in enumerate(col_widths):
        table.columns[col].width = Inches(width)
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = BLUE
    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = DARK
    return slide


def add_callout(slide, text, left=0.55, top=6.1, width=12.0, height=0.5):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = MID_BLUE
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.LEFT


def main() -> None:
    milestones = load_schedule_summary()
    top_risks = load_risk_summary()
    labour_total, capex_total, categories = load_budget_summary()

    prs = Presentation(str(TEMPLATE_PATH)) if TEMPLATE_PATH.exists() else Presentation()
    clear_slides(prs)

    cover = prs.slides.add_slide(prs.slide_layouts[6])
    cover.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5)).fill.solid()
    cover.shapes[-1].fill.fore_color.rgb = BLUE
    cover.shapes[-1].line.color.rgb = BLUE
    cover.shapes.add_picture(str(LOGO_PATH), Inches(0.7), Inches(0.6), height=Inches(0.55))
    title_box = cover.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.5), Inches(1.0))
    p = title_box.text_frame.paragraphs[0]
    p.text = PROJECT_NAME
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    sub = title_box.text_frame.add_paragraph()
    sub.text = f"{BUSINESS} | {MODEL} model | Stakeholder presentation"
    sub.font.size = Pt(18)
    sub.font.color.rgb = RGBColor(255, 255, 255)
    meta = cover.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(5.8), Inches(1.6))
    mp = meta.text_frame.paragraphs[0]
    mp.text = f"Seller: {SELLER}"
    mp.font.size = Pt(18)
    mp.font.color.rgb = RGBColor(255, 255, 255)
    for line in [f"Buyer / JV partner: {BUYER}", f"Report date: {REPORT_DATE.strftime('%d %B %Y')}", "Baseline version: v1.0"]:
        para = meta.text_frame.add_paragraph()
        para.text = line
        para.font.size = Pt(16)
        para.font.color.rgb = RGBColor(255, 255, 255)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_branding(slide, "Executive Summary", "Management view from the current Gamma baseline outputs")
    add_bullets(
        slide,
        [
            "Gamma separates the Bosch Cloud business from Robert Bosch China and establishes the Day 1 operating baseline for a 50/50 JV with Alibaba.",
            "Scope is intentionally lean for Day 1: 5 sites, 250 users, 20 applications, and no SAP. Delivery concentration is therefore on infrastructure, identity, operations, and confidentiality control.",
            "The current status is a pre-start baseline: the schedule, risk register, cost plan, charter, dashboards, and monthly report are complete, but execution is still dependent on legal and disclosure timing.",
            "Executive focus remains on antitrust due diligence, clean-team governance, and early confirmation of shared-service dependencies before mobilization starts on 01 Aug 2026.",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_branding(slide, "Scope And Operating Model", "What Gamma must deliver at Day 1")
    add_bullets(
        slide,
        [
            "In scope: infrastructure separation, identity and security redesign, service management readiness, transition of 20 applications, and operational handover for 250 users across 5 sites.",
            "Out of scope: SAP, manufacturing / OT, and broad Day 2 transformation beyond the minimum viable JV operating baseline.",
            "Combination model: the project may use interim service continuity during transition, but the target state is a jointly managed JV environment rather than a permanent shared platform.",
            "The 50/50 governance structure requires clear decision rights before QG1 and again before QG4, because ambiguity at either gate would directly erode GoLive confidence.",
        ],
    )

    timeline_slide = add_table_slide(
        prs,
        "Timeline And Quality Gates",
        "Milestones reconciled to the generated Gamma schedule",
        ["Milestone", "Date"],
        milestones,
        [8.5, 3.2],
    )
    milestone_lookup = dict(milestones)
    add_callout(
        timeline_slide,
        f"Key dates: QG1 {milestone_lookup['QG1 - Concept approved']} | GoLive {milestone_lookup['GoLive - Day 1 cutover complete']} | QG5 {milestone_lookup['QG5 - Project completion approved']}",
    )

    budget_rows = [(name, f"EUR {value:,.0f}") for name, value in categories]
    budget_rows.append(("Labour baseline total", f"EUR {labour_total:,.0f}"))
    budget_rows.append(("CAPEX / additional costs", f"EUR {capex_total:,.0f}"))
    budget_slide = add_table_slide(
        prs,
        "Budget And Cost Structure",
        "Labour baseline separated from risk-driven contingency",
        ["Category", "Amount"],
        budget_rows,
        [8.5, 3.2],
    )
    add_callout(budget_slide, f"Budget anchor: Labour EUR {labour_total:,.0f} | CAPEX / additional EUR {capex_total:,.0f} | Draft baseline for QG1 approval")

    risk_rows = [(f"#{risk_id}", category, score, event) for risk_id, score, category, event in top_risks]
    risk_slide = add_table_slide(
        prs,
        "Top Risks And Management Actions",
        "Threat risks reconciled to the Gamma risk register",
        ["Risk", "Category", "Score", "Executive issue"],
        risk_rows,
        [0.9, 2.4, 0.8, 8.1],
    )
    add_callout(
        risk_slide,
        "Risk anchors: " + " | ".join(f"#{risk_id} {category} ({score})" for risk_id, score, category, _ in top_risks[:5]) + " | #16 Schedule (15)",
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_branding(slide, "Decisions Required", "Immediate steerco focus for a clean mobilisation")
    add_bullets(
        slide,
        [
            "Confirm the legal and antitrust decision path that will govern clean-team expansion and future stakeholder disclosure.",
            "Approve the named clean-team resource pool across PMO, legal, infrastructure, security, and operations before the 01 Aug 2026 kickoff.",
            "Approve the draft labour baseline and associated contingency reserve for QG1 so that external legal, provider, and readiness support can be activated quickly if needed.",
            "Endorse the principle that Day 1 scope stays intentionally lean - no SAP, no broad Day 2 transformation, and no migration activity after GoLive.",
        ],
    )

    prs.save(OUTPUT_PATH)
    print(f"[{PROJECT_NAME}] Stakeholder presentation written to {OUTPUT_PATH}")
    print(f"  Slides generated: {len(prs.slides)}")


if __name__ == "__main__":
    main()