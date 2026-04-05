from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# Bosch palette used across existing deliverables in this repository.
BOSCH_NAVY = RGBColor(0x00, 0x3B, 0x6E)
BOSCH_BLUE = RGBColor(0x00, 0x66, 0xCC)
BOSCH_LIGHT = RGBColor(0xE4, 0xED, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5A, 0x64, 0x78)
GREEN = RGBColor(0x00, 0x7A, 0x33)
AMBER = RGBColor(0xE8, 0xA0, 0x00)
RED = RGBColor(0xCC, 0x00, 0x00)

ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "Charlie" / "Charlie_Management_Stakeholder_Presentation.pptx"
PPT_TEMPLATE = ROOT / "Bosch presentation template.pptx"


def find_bosch_logo() -> Path:
    candidates = [
        ROOT / "Bosch.png",
        ROOT / "Reference" / "Bosch.png",
        ROOT / "Licensing" / "Bosch.png",
    ]
    for candidate in candidates:
        if candidate.exists():
            return candidate
    raise FileNotFoundError("Bosch logo not found in expected locations.")


def add_header(slide, title: str, subtitle: str, logo_path: Path) -> None:
    bg = slide.shapes.add_shape(1, 0, 0, Inches(13.33), Inches(1.25))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BOSCH_NAVY
    bg.line.fill.background()

    slide.shapes.add_picture(str(logo_path), Inches(0.35), Inches(0.2), height=Inches(0.45))

    title_box = slide.shapes.add_textbox(Inches(1.8), Inches(0.2), Inches(8.5), Inches(0.45))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE

    sub_box = slide.shapes.add_textbox(Inches(1.8), Inches(0.7), Inches(10.5), Inches(0.35))
    p2 = sub_box.text_frame.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(12)
    p2.font.color.rgb = WHITE


def add_footer(slide, text: str) -> None:
    footer = slide.shapes.add_textbox(Inches(0.4), Inches(7.1), Inches(12.4), Inches(0.3))
    p = footer.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def add_bullets(slide, x, y, w, h, lines):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18 if idx == 0 else 15)
        p.font.bold = idx == 0
        p.font.color.rgb = TEXT_DARK
        if idx > 0:
            p.space_after = Pt(10)


def add_card(slide, x, y, w, h, title, value, color=BOSCH_LIGHT):
    card = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.12), Inches(w - 0.3), Inches(0.35))
    t = title_box.text_frame.paragraphs[0]
    t.text = title
    t.font.size = Pt(11)
    t.font.bold = True
    t.font.color.rgb = MUTED

    value_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.55), Inches(w - 0.3), Inches(h - 0.7))
    v = value_box.text_frame.paragraphs[0]
    v.text = value
    v.font.size = Pt(20)
    v.font.bold = True
    v.font.color.rgb = BOSCH_NAVY


def create_presentation() -> Presentation:
    if PPT_TEMPLATE.exists():
        return Presentation(str(PPT_TEMPLATE))

    print(
        f"Warning: template not found at {PPT_TEMPLATE}. "
        "Falling back to default presentation theme."
    )
    return Presentation()


def set_placeholder_text(slide, idx, text) -> bool:
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            shape.text = text
            return True
    return False


def add_cover_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_placeholder_text(slide, 0, "Project Charlie")
    subtitle = (
        "Management Stakeholder Update\n"
        "Robert Bosch GmbH AI Business Carve-Out into 50/50 JV with Undisclosed\n"
        "Report Date: 04 Apr 2026"
    )
    if not set_placeholder_text(slide, 1, subtitle):
        box = slide.shapes.add_textbox(Inches(1.0), Inches(4.9), Inches(10.5), Inches(1.2))
        box.text_frame.text = subtitle


def add_timeline_slide(prs: Presentation, title: str, subtitle: str) -> None:
    """Special handling for Timeline and Quality Gates slide with visual timeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[10])

    if not set_placeholder_text(slide, 0, title):
        t = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(10.0), Inches(0.5))
        t.text_frame.text = title

    if not set_placeholder_text(slide, 15, subtitle):
        st = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.0), Inches(0.5))
        st.text_frame.text = subtitle

    # Timeline data: (label, date, x_position, color)
    milestones = [
        ("QG0", "30 Apr", 1.0, BOSCH_BLUE),
        ("QG1/2/3", "26 Jun", 4.2, RGBColor(0x00, 0x88, 0xE0)),
        ("GoLive", "01 Jul", 5.5, GREEN),
        ("QG4", "02 Jul", 5.8, RGBColor(0x00, 0x76, 0xC9)),
        ("QG5", "30 Oct", 10.2, BOSCH_LIGHT),
    ]

    # Draw connecting timeline line
    line = slide.shapes.add_shape(1, Inches(0.8), Inches(2.9), Inches(11.0), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = BOSCH_LIGHT
    line.line.fill.background()

    # Draw milestones
    for label, date, x_pos, color in milestones:
        # Milestone circle
        circle = slide.shapes.add_shape(9, Inches(x_pos - 0.2), Inches(2.65), Inches(0.4), Inches(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color if color != BOSCH_LIGHT else BOSCH_BLUE
        circle.line.color.rgb = WHITE
        circle.line.width = Pt(2)

        # Label box below
        label_box = slide.shapes.add_textbox(Inches(x_pos - 0.4), Inches(3.2), Inches(0.8), Inches(0.35))
        lp = label_box.text_frame.paragraphs[0]
        lp.text = label
        lp.alignment = PP_ALIGN.CENTER
        lp.font.size = Pt(12)
        lp.font.bold = True
        lp.font.color.rgb = BOSCH_NAVY

        # Date below label
        date_box = slide.shapes.add_textbox(Inches(x_pos - 0.4), Inches(3.55), Inches(0.8), Inches(0.3))
        dp = date_box.text_frame.paragraphs[0]
        dp.text = date
        dp.alignment = PP_ALIGN.CENTER
        dp.font.size = Pt(10)
        dp.font.color.rgb = MUTED

    # Add key notes below timeline
    notes_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.5), Inches(1.2))
    tf = notes_box.text_frame
    tf.clear()
    
    notes = [
        "Key milestones: QG0 (30 Apr) â†’ QG1/2/3 combined (26 Jun) â†’ GoLive (01 Jul) and QG4 (02 Jul) â†’ QG5 (30 Oct)",
        "QG4 Pre-GoLive gate must be cleared before Day 1 operations commence on 01 Jun 2027",
    ]
    for i, note in enumerate(notes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = note
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(6)

    # Add top accent bar only (no graphics/chips for timeline)
    top_bar = slide.shapes.add_shape(1, Inches(0.22), Inches(1.33), Inches(11.55), Inches(0.06))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = BOSCH_BLUE
    top_bar.line.fill.background()


def add_content_slide(prs: Presentation, title: str, subtitle: str, bullets) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[10])

    if not set_placeholder_text(slide, 0, title):
        t = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(10.0), Inches(0.5))
        t.text_frame.text = title

    if not set_placeholder_text(slide, 15, subtitle):
        st = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.0), Inches(0.5))
        st.text_frame.text = subtitle

    content_placeholder = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            content_placeholder = shape
            break

    if content_placeholder is None:
        content_placeholder = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.8))

    tf = content_placeholder.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(22 if i == 0 else 18)
        p.font.bold = i == 0
        p.font.color.rgb = TEXT_DARK

    add_professional_graphics(slide, title)


def get_slide_chips(title: str):
    chips = {
        "Executive Summary": [("GoLive", "01 Jun 2027"), ("SPI", "1.00"), ("CPI", "1.00")],
        "Scope and Carve-Out Model": [("Applications", "17"), ("Users", "70"), ("Sites", "2")],
        "Timeline and Quality Gates": [("QG0", "30 Apr"), ("QG4", "02 Jul"), ("QG5", "30 Oct")],
        "KPI and Readiness Snapshot": [("Open Risks", "3 Amber"), ("Critical Issues", "0"), ("Readiness", "On Plan")],
        "Budget View": [("Labour", "EUR 554K"), ("CAPEX", "TBC QG0"), ("Reserve", "EUR 83K")],
        "Top Risks and Mitigations": [("R1", "16"), ("R5", "12"), ("Review", "Weekly")],
        "Management Decisions Required": [("Decisions", "4"), ("Owner", "SteerCo"), ("Cadence", "Weekly")],
    }
    return chips.get(title, [("Status", "Active"), ("Plan", "On Track"), ("Review", "Weekly")])


def add_professional_graphics(slide, title: str) -> None:
    # Thin accent separator below title/subtitle area.
    top_bar = slide.shapes.add_shape(1, Inches(0.22), Inches(1.33), Inches(11.55), Inches(0.06))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = BOSCH_BLUE
    top_bar.line.fill.background()

    # Subtle accent motif in the top-right to avoid footer overlap.
    c3 = slide.shapes.add_shape(9, Inches(11.65), Inches(0.25), Inches(0.22), Inches(0.22))
    c3.fill.solid()
    c3.fill.fore_color.rgb = BOSCH_BLUE
    c3.line.fill.background()

    # Add three compact chips in the middle area for quick executive scan.
    chips = get_slide_chips(title)
    x = 0.3
    for label, value in chips:
        chip = slide.shapes.add_shape(1, Inches(x), Inches(4.9), Inches(3.65), Inches(0.58))
        chip.fill.solid()
        chip.fill.fore_color.rgb = BOSCH_LIGHT
        chip.line.fill.background()

        label_box = slide.shapes.add_textbox(Inches(x + 0.12), Inches(4.96), Inches(1.8), Inches(0.18))
        lp = label_box.text_frame.paragraphs[0]
        lp.text = label
        lp.font.size = Pt(10)
        lp.font.bold = True
        lp.font.color.rgb = MUTED

        value_box = slide.shapes.add_textbox(Inches(x + 0.12), Inches(5.14), Inches(3.35), Inches(0.24))
        vp = value_box.text_frame.paragraphs[0]
        vp.text = value
        vp.font.size = Pt(13)
        vp.font.bold = True
        vp.font.color.rgb = BOSCH_NAVY
        x += 3.84


def clear_all_slides(prs: Presentation) -> None:
    # The Bosch template can include example slides; remove them before generating.
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]


def main():
    prs = create_presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    clear_all_slides(prs)

    add_cover_slide(prs)

    slides = [
        (
            "Executive Summary",
            "Project Charlie | Stand Alone model with stand-alone separation",
            [
                "Status: On track for Day 1 GoLive on 01 Jun 2027",
                "Scope: 17 AI applications, 3500+ users, 37 worldwide sites.",
                "Model benefits: no TSA, no Merger Zone, no antitrust filing.",
                "Baseline: EUR 554K labour budget, CAPEX decision at QG0.",
                "SteerCo focus: legal entity timing and Jun build resourcing.",
            ],
        ),
        (
            "Scope and Carve-Out Model",
            "What is in scope and why execution risk is lower",
            [
                "In scope: AI app wave, users/devices, cloud/identity/network/security",
                "17 AI applications in one migration wave",
                "3500+ users and device transition to JV standards",
                "JV AD forest, M365 tenant, Azure setup across 37 worldwide sites",
                "No ERP in scope and no TSA required after GoLive",
            ],
        ),
        (
            "Timeline and Quality Gates",
            "April to October 2026 plan with mandatory gate control",
            [],
        ),
        (
            "KPI and Readiness Snapshot",
            "Control metrics and confidence by workstream",
            [
                "SPI 1.00 and CPI 1.00: on plan",
                "Apps migrated: 0/17 (migration planned in June)",
                "Open risks: 3 amber | Critical issues: 0",
                "High confidence: PMO, Infra, Security, Workplace",
                "Medium confidence: Legal setup, Data separation, Licensing",
            ],
        ),
        (
            "Budget View",
            "Labour baseline and major cost buckets",
            [
                "Labour baseline: EUR 554K | CAPEX: TBC at QG0",
                "Top cost buckets: Programme Mgmt, IT PM, Hypercare",
                "Infrastructure and app migration costs within plan",
                "Contingency reserve available for transition volatility",
                "No budget overrun indicators at current stage",
            ],
        ),
        (
            "Top Risks and Mitigations",
            "Active amber risks and mitigation status",
            [
                "R1 Bandwidth constraints Apr-Jun (score 16, amber)",
                "R3 Licensing cost pressure (score 9, amber)",
                "R5 Key IT staff availability in June (score 12, amber)",
                "Mitigations active: resource protection, quote lock, staffing backup",
                "Weekly escalation in PMO governance",
            ],
        ),
        (
            "Management Decisions Required",
            "Steering actions requested this month",
            [
                "Approve CAPEX envelope at QG0",
                "Confirm Undisclosed-side IT staffing for May architecture workshops",
                "Approve June critical resource protection plan",
                "Support legal escalation for JV entity registration timeline",
                "Maintain weekly executive checkpoint cadence to GoLive",
            ],
        ),
    ]

    for idx, (title, subtitle, bullets) in enumerate(slides):
        if title == "Timeline and Quality Gates":
            add_timeline_slide(prs, title, subtitle)
        else:
            add_content_slide(prs, title, subtitle, bullets)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Created: {OUTPUT}")


if __name__ == "__main__":
    main()


