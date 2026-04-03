#!/usr/bin/env python3
"""
Generate Project Trinity Project Charter PPTX.

Strategy
--------
1. Copy the Bosch template PPTX.
2. Replace every '<name of project>' / '<activity Name>' / '<date>' placeholder
   in all text runs with the real project name / date.
3. For each content slide, render a styled image (matplotlib → PNG) and paste it
   at the content-area position, so fonts, alignment and colour are under full
   control.  The slide title, breadcrumb header and footer are kept from the
   template (only their text is updated).
"""

import argparse
import io
import shutil
import textwrap
from pathlib import Path

import warnings
warnings.filterwarnings("ignore", message=".*tight_layout.*")
warnings.filterwarnings("ignore", message=".*Axes.*not compatible.*")
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.table import Table
import numpy as np
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# ── Paths ─────────────────────────────────────────────────────────────────────
BASE_DIR = Path(__file__).parent
TEMPLATE  = BASE_DIR / "Project_charter_template.pptx"

# ── Design tokens ─────────────────────────────────────────────────────────────
BLUE   = "#0070C0"    # Bosch / template primary blue
DKGRAY = "#333333"
LTBLUE = "#DDEEFF"    # light blue for alternate table rows
WHITE  = "#FFFFFF"
BLACK  = "#000000"
RED    = "#C00000"    # for risk HIGH cells
AMBER  = "#E08000"
GREEN  = "#00703C"

FONT   = "Arial"
DPI    = 150          # image resolution

# ── Project Data ──────────────────────────────────────────────────────────────
P = {
    # ── Identity ──────────────────────────────────────────────────────────────
    # Project FRAME (Bosch → Keenfinity) is the REFERENCE document only.
    # Project Trinity is a NEW engagement: Bosch carves out a business to JCI.
    "name":             "Project Trinity",
    # JCI (seller) divests heating business → Bosch (buyer) acquires and integrates
    "full":             "IT Carve-Out: JCI Heating Business to Bosch",
    "dept":             "Bosch IT / JCI IT / KPMG",
    "pm":               "Bosch IT PM + JCI IT PM",
    "seller":           "JCI (Johnson Controls International)",
    "buyer":            "Bosch",
    "sponsor_customer": "Bosch  (Buyer — acquiring JCI Heating Business)",
    "sponsor_contr":    "JCI  (Seller — IT operations during TSA)",
    "sc_members":       "Bosch Executive Leadership + JCI Executive Leadership",
    "delivery":         "Bosch IT Team, JCI IT Team, KPMG (Methodology Lead)",
    "supporting":       "Legal, Procurement, Regional IT Leads (AP / AM / EMEA)",
    # ── Dates ─────────────────────────────────────────────────────────────────
    "date":             "29.03.2026",
    "type":             "Major IT Programme",
    "kickoff":          "01.07.2026",
    "qg1":              "30.09.2026",
    "qg2":              "31.03.2027",
    "qg3":              "31.07.2027",
    "qg4":              "31.10.2027",
    "day1":             "01.11.2027",
    "closure":          "28.02.2028",
    # ── Scope ─────────────────────────────────────────────────────────────────
    "employees":        "8,000",
    "sites":            "180",
    "apps":             "~500",
    "devices":          "~6,000",
    "regions":          "AP / AM / EMEA",
    # ── Budget ────────────────────────────────────────────────────────────────
    "budget_labour":    "EUR 5.1M",
    "budget_gov":       "EUR 633.6K",
    "budget_bosch":     "EUR 1.95M",
    "budget_jci":       "EUR 1.2M",
    "budget_kpmg":      "EUR 1.3M",
    "ftes_peak":        "~45 FTE",
    "tsa":              "18 months",
    "model":            "Stand Alone (Full Independence)",
    "duration":         "24 months (Jul 2026 - Feb 2028)",
    "charter_title":    "Project Trinity - IT Carve-Out: Bosch to JCI",
}

# Text fragments to swap in all runs (avoids <> / & XML issues)
GLOBAL_SWAPS = [
    ("(Name of project)",                        P["name"]),
    ("<name of project>",                        P["name"]),
    ("< name of project >",                      P["name"]),
    ("<name of project",                         P["name"]),
    ("<activity Name>",                          P["name"]),
    ("<Date>",                                   f"<{P['date']}>"),
    ("Project Type: Carve out",                  f"Project Type: {P['type']}"),
    ("Huang, Christina CI/MIR-AP",               P["pm"]),
    ("CI/MIO",                                   "Bosch IT / KPMG"),
]

# ── Helpers ───────────────────────────────────────────────────────────────────

def replace_globally(prs: Presentation) -> None:
    """Swap placeholder text in every run of every slide."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in GLOBAL_SWAPS:
                        if old in run.text:
                            run.text = run.text.replace(old, new)


def fig_to_pil(fig) -> Image.Image:
    """Convert a matplotlib Figure to a PIL Image."""
    buf = io.BytesIO()
    with warnings.catch_warnings():
        warnings.simplefilter("ignore")
        fig.savefig(buf, format="png", dpi=DPI, bbox_inches="tight",
                    facecolor=fig.get_facecolor())
    buf.seek(0)
    img = Image.open(buf).convert("RGBA")
    return img


def pil_to_stream(img: Image.Image) -> io.BytesIO:
    stream = io.BytesIO()
    img.save(stream, format="PNG")
    stream.seek(0)
    return stream


def insert_image(slide, img: Image.Image,
                 left_in: float, top_in: float,
                 width_in: float, height_in: float) -> None:
    """Insert a PIL image onto a slide at given position (inches)."""
    stream = pil_to_stream(img)
    slide.shapes.add_picture(
        stream,
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in),
    )


def hide_placeholder(slide, ph_idx: int = 1) -> tuple[float, float, float, float]:
    """
    Clear text in body placeholder (idx=ph_idx) and return its position.
    Returns (left_in, top_in, width_in, height_in).
    """
    EMU = 914400
    for shape in slide.shapes:
        try:
            ph = shape.placeholder_format
            if ph and ph.idx == ph_idx:
                # Clear text
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ""
                # Also remove bullet symbols / non-run text
                for p_elem in shape.text_frame._txBody.findall(qn("a:p")):
                    for r in p_elem.findall(qn("a:r")):
                        for t in r.findall(qn("a:t")):
                            t.text = ""
                return (shape.left / EMU, shape.top / EMU,
                        shape.width / EMU, shape.height / EMU)
        except Exception:
            pass
    # fallback: standard content area
    return (0.283, 1.25, 11.43, 4.79)


def make_fig(w_in: float, h_in: float) -> plt.Figure:
    """Return a blank white figure of given size in inches."""
    fig = plt.figure(figsize=(w_in, h_in), facecolor=WHITE)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis("off")
    return fig, ax


# ── Low-level drawing primitives ──────────────────────────────────────────────

def draw_section_header(ax, y: float, label: str, x0: float = 0.01,
                        fontsize: float = 11) -> None:
    """Blue bold section label with a thin underline."""
    ax.text(x0, y, label, color=BLUE, fontsize=fontsize, fontfamily=FONT,
            fontweight="bold", va="top", transform=ax.transAxes)
    # Draw underline using plot in axes-fraction coordinates
    ax.plot([x0, 0.99], [y - 0.012, y - 0.012], color=BLUE, linewidth=0.8,
            transform=ax.transAxes, clip_on=False)


def draw_bullet_block(ax, y_start: float, items: list[str],
                      x0: float = 0.03, fontsize: float = 9.5,
                      line_gap: float = 0.042,
                      max_chars: int = 85) -> float:
    """Draw a bulleted text block with wrapping; returns y after last line."""
    y = y_start
    for item in items:
        is_sub = item.startswith("  ")
        bullet = " " if is_sub else "•"
        text = f"{bullet} {item.lstrip()}"
        wrapped = textwrap.wrap(text, width=max_chars) or [text]
        for line in wrapped:
            ax.text(x0, y, line, color=DKGRAY, fontsize=fontsize,
                    fontfamily=FONT, va="top", transform=ax.transAxes,
                    clip_on=True)
            y -= line_gap
    return y


def draw_table(ax, rows: list[list[str]], col_widths: list[float],
               header: bool = True, x0: float = 0.01, y0: float = 0.97,
               row_height: float = 0.07, fontsize: float = 8.5,
               header_color: str = BLUE, alt_color: str = LTBLUE) -> None:
    """
    Draw a styled table using matplotlib patches + text.
    rows[0] is treated as the header row when header=True.
    col_widths: fractions of axes width, must sum <= 1.
    """
    # Build x positions from col_widths
    x_positions = [x0]
    for w in col_widths[:-1]:
        x_positions.append(x_positions[-1] + w)
    total_w = sum(col_widths)

    for r_idx, row in enumerate(rows):
        y_top = y0 - r_idx * row_height
        y_bot = y_top - row_height
        is_header = header and r_idx == 0
        bg = header_color if is_header else (alt_color if r_idx % 2 == 0 else WHITE)
        txt_color = WHITE if is_header else DKGRAY
        fw = "bold" if is_header else "normal"

        # Draw row background
        rect = mpatches.FancyBboxPatch(
            (x0, y_bot), total_w, row_height,
            boxstyle="square,pad=0", linewidth=0.3,
            edgecolor="#AAAAAA", facecolor=bg,
            transform=ax.transAxes, clip_on=True,
        )
        ax.add_patch(rect)

        for c_idx, (cell, cx, cw) in enumerate(zip(row, x_positions, col_widths)):
            ax.text(cx + 0.008, y_bot + row_height * 0.5, str(cell),
                    color=txt_color, fontsize=fontsize, fontfamily=FONT,
                    fontweight=fw, va="center", transform=ax.transAxes,
                    clip_on=True)


# ── Per-slide image generators ────────────────────────────────────────────────

def img_agenda(w, h) -> Image.Image:
    fig, ax = make_fig(w, h)
    items = [
        ("1.", "Roles & Responsibilities"),
        ("2.", "Goals & Success Criteria"),
        ("3.", "Scope Description"),
        ("4.", "Solution Architecture (Carve-Out Model)"),
        ("5.", "Business Case"),
        ("6.", "Governance"),
        ("7.", "Assumptions & Constraints"),
        ("8.", "Development / Delivery Approach"),
        ("9.", "Risk Register (Top 5)"),
        ("10.", "Timeline & Milestones"),
        ("11.", "Cost / Budget"),
        ("12.", "Resources"),
        ("13.", "Project Organisation"),
        ("14.", "Communication Plan"),
        ("15.", "Status Report"),
    ]
    y = 0.95
    for num, label in items:
        ax.text(0.03, y, num, color=BLUE, fontsize=11, fontfamily=FONT,
                fontweight="bold", va="top", transform=ax.transAxes)
        ax.text(0.10, y, label, color=DKGRAY, fontsize=11, fontfamily=FONT,
                va="top", transform=ax.transAxes)
        ax.plot([0.02, 0.98], [y - 0.038, y - 0.038], color="#DDDDDD",
                linewidth=0.5, transform=ax.transAxes, clip_on=False)
        y -= 0.058
    plt.tight_layout(pad=0.1)
    img = fig_to_pil(fig)
    plt.close(fig)
    return img


def img_roles(w, h) -> Image.Image:
    fig, ax = make_fig(w, h)
    rows = [
        ["Role", "Name / Organisation", "R", "A", "C", "I"],
        ["Project Sponsor", "Bosch CIO / JCI EVP", "", "A", "", ""],
        ["IT Programme Manager", "Bosch IT PM", "R", "A", "", ""],
        ["IT Project Manager", "JCI IT PM", "R", "A", "", ""],
        ["Consulting Lead", "KPMG Engagement Lead", "R", "", "C", ""],
        ["Steering Committee", "Bosch CIO, JCI EVP, CFOs, Legal", "", "A", "", "I"],
        ["WS Lead - Infrastructure", "Bosch IT Infrastructure", "R", "", "C", ""],
        ["WS Lead - ERP / SAP", "JCI ERP Team + KPMG", "R", "", "C", ""],
        ["WS Lead - Applications", "JCI IT Architects", "R", "", "C", ""],
        ["WS Lead - HR IT", "JCI HR IT", "R", "", "C", ""],
        ["WS Lead - IT Security", "Bosch CISO", "R", "", "C", ""],
        ["Legal / Contracts", "Bosch Legal + JCI Legal", "", "", "C", "I"],
        ["Regional Leads", "AP / AM / EMEA IT Leads", "R", "", "C", ""],
    ]
    # col widths: role, name, R, A, C, I
    cw = [0.30, 0.45, 0.07, 0.07, 0.07, 0.07]
    draw_table(ax, rows, cw, y0=0.97, row_height=0.066, fontsize=9)

    # RACI legend
    ax.text(0.01, 0.08, "R = Responsible   A = Accountable   C = Consulted   I = Informed",
            color=DKGRAY, fontsize=8, fontfamily=FONT, transform=ax.transAxes)
    plt.tight_layout(pad=0.1)
    img = fig_to_pil(fig)
    plt.close(fig)
    return img


def img_goals(w, h) -> Image.Image:
    """Target Cross — 4-quadrant project goals layout (Bosch to JCI).

    Source data: Trinity_Project_Schedule.csv, Trinity_Status_Dashboard.html,
    Trinity_Resource_Summary.csv, Trinity_Project_Cost_Plan.csv.
    Project FRAME (Bosch to Keenfinity) is the REFERENCE document only;
    no FRAME-specific content appears here.
    """
    import textwrap

    # ── Content per quadrant (sentences, not bullet points) ───────────────────
    # Business flow: JCI (seller) divests heating business to Bosch (buyer).
    # IT currently runs on JCI infrastructure. Carve-out separates the heating
    # business IT from JCI, stages it in a Merger Zone, then integrates it
    # into Bosch's IT environment. JCI operates IT under 18-month TSA while
    # Bosch builds the receiving infrastructure.
    QUADRANTS = [
        {
            "label": "Business Needs",
            "pos":   "top-left",
            "color": "#004B87",
            "text": (
                "Bosch is acquiring JCI's heating business and must carve out all "
                "associated IT systems, applications and data from JCI's infrastructure "
                "and integrate them into the Bosch IT environment. "
                "The heating business spans 8,000 employees across 180 sites in AP, AM "
                "and EMEA, with ~500 applications and 6,000 client devices currently "
                "running on JCI's IT landscape. "
                "JCI will continue operating IT for the carved-out business under an "
                "18-month TSA while Bosch builds the receiving infrastructure to fully "
                "absorb and integrate the heating business."
            ),
        },
        {
            "label": "Sponsor / Stakeholder",
            "pos":   "top-right",
            "color": "#0070C0",
            "lines": [
                ("Sponsor Customer",   "Bosch  (Buyer — acquiring JCI Heating Business)"),
                ("Sponsor Contractor", "JCI  (Seller — IT operations during TSA)"),
                ("", ""),
                ("Steering Committee", "Bosch Executive Leadership + JCI Executive Leadership"),
                ("Delivery Partners",  "Bosch IT Team, JCI IT Team, KPMG (Methodology Lead)"),
                ("Supporting",         "Legal, Procurement, Regional IT Leads (AP / AM / EMEA)"),
            ],
        },
        {
            "label": "Result",
            "pos":   "bottom-left",
            "color": "#0095C8",
            "text": (
                "The JCI heating business IT is fully carved out from JCI's infrastructure "
                "and integrated into the Bosch IT environment. "
                "All 8,000 employees, ~500 applications, SAP landscape, M365 tenant and "
                "identity management are migrated to Bosch systems via a Merger Zone "
                "interim environment, with no remaining dependency on JCI's IT. "
                "The 18-month TSA is successfully exited, and the heating business "
                "operates entirely on Bosch infrastructure."
            ),
        },
        {
            "label": "Criteria of Success",
            "pos":   "bottom-right",
            "color": "#00A3AD",
            "text": (
                f"All 8,000 heating business employees are fully operational on Bosch IT "
                f"infrastructure on Day 1 ({P['day1']}) with zero unplanned downtime "
                f"on ERP, email and network systems. "
                f"At least 95% of applications are migrated or decommissioned on schedule "
                f"and all JCI contracts and licences are transferred to Bosch with no "
                f"legal blockers. "
                f"The 18-month TSA with JCI is exited with fewer than 2% escalation rate "
                f"during the 3-month hypercare period, and the project is delivered within "
                f"the {P['budget_labour']} labour budget with all six quality gates passed."
            ),
        },
    ]

    GAP    = 0.012   # cross gap between quadrants (axes fraction)
    HDR_H  = 0.115   # header band height (axes fraction)
    PAD    = 0.025   # text padding inside each box

    fig = plt.figure(figsize=(w, h), facecolor=WHITE)
    ax  = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis("off")

    # ── Quadrant boundaries ───────────────────────────────────────────────────
    # x: [0 .. mid-GAP/2] and [mid+GAP/2 .. 1]
    # y: [0 .. mid-GAP/2] and [mid+GAP/2 .. 1]
    mid = 0.5
    boxes = {
        "top-left":     (0,            mid + GAP/2, mid - GAP/2, 1),
        "top-right":    (mid + GAP/2,  mid + GAP/2, 1,           1),
        "bottom-left":  (0,            0,            mid - GAP/2, mid - GAP/2),
        "bottom-right": (mid + GAP/2,  0,            1,           mid - GAP/2),
    }

    # ── Small center diamond ──────────────────────────────────────────────────
    diamond_verts = [
        (mid, mid - 0.038),
        (mid + 0.028, mid),
        (mid, mid + 0.038),
        (mid - 0.028, mid),
    ]
    diamond = plt.Polygon(diamond_verts, closed=True,
                          facecolor=WHITE, edgecolor="#CCCCCC",
                          linewidth=1.0, transform=ax.transAxes, zorder=5)
    ax.add_patch(diamond)

    # ── Draw each quadrant ────────────────────────────────────────────────────
    for q in QUADRANTS:
        x0, y0, x1, y1 = boxes[q["pos"]]
        w_box  = x1 - x0
        h_box  = y1 - y0
        color  = q["color"]

        # Background
        bg = mpatches.FancyBboxPatch(
            (x0, y0), w_box, h_box,
            boxstyle="square,pad=0",
            facecolor="#F0F7FF", edgecolor="#CCCCCC", linewidth=0.8,
            transform=ax.transAxes, zorder=1,
        )
        ax.add_patch(bg)

        # Header band
        hdr_y0 = y1 - HDR_H
        hdr = mpatches.FancyBboxPatch(
            (x0, hdr_y0), w_box, HDR_H,
            boxstyle="square,pad=0",
            facecolor=color, edgecolor="none",
            transform=ax.transAxes, zorder=2,
        )
        ax.add_patch(hdr)

        # Header label (centred in band)
        ax.text(
            x0 + w_box / 2, hdr_y0 + HDR_H / 2,
            q["label"],
            color=WHITE, fontsize=11, fontfamily=FONT, fontweight="bold",
            ha="center", va="center", transform=ax.transAxes, zorder=3,
        )

        # ── Body text ─────────────────────────────────────────────────────────
        # available body area
        body_y_top = hdr_y0 - PAD          # just below header
        body_h     = hdr_y0 - y0 - 2 * PAD
        body_x0    = x0 + PAD
        body_w     = w_box - 2 * PAD

        # Estimate chars per line from body width (rough: ~1 char ≈ 0.007 axes units at fontsize 9)
        chars_per_line = max(20, int(body_w / 0.0062))

        if "text" in q:
            # Flowing paragraph — wrap and render line by line
            import textwrap as tw
            wrapped = tw.fill(q["text"], width=chars_per_line)
            lines   = wrapped.split("\n")
            line_h  = body_h / max(len(lines) + 1, 1)
            fontsize = min(9.5, max(7.5, body_h * 72 / (len(lines) + 2)))

            for i, line in enumerate(lines):
                y_pos = body_y_top - i * (body_h / len(lines)) - 0.01
                ax.text(
                    body_x0, y_pos, line,
                    color=DKGRAY, fontsize=fontsize, fontfamily=FONT,
                    va="top", ha="left", transform=ax.transAxes, zorder=3,
                )

        elif "lines" in q:
            # Structured label: value pairs
            n = sum(1 for lbl, val in q["lines"] if lbl or val)
            fontsize = min(9.5, max(7.5, body_h * 72 / (n + 2)))
            line_step = body_h / max(n + 1, 1)
            y_cursor  = body_y_top - 0.01
            for lbl, val in q["lines"]:
                if not lbl and not val:
                    y_cursor -= line_step * 0.4
                    continue
                ax.text(
                    body_x0, y_cursor, lbl,
                    color=color, fontsize=fontsize, fontfamily=FONT,
                    fontweight="bold", va="top", ha="left",
                    transform=ax.transAxes, zorder=3,
                )
                ax.text(
                    body_x0 + 0.13, y_cursor, val,
                    color=DKGRAY, fontsize=fontsize, fontfamily=FONT,
                    va="top", ha="left", transform=ax.transAxes, zorder=3,
                )
                y_cursor -= line_step

    plt.tight_layout(pad=0)
    img = fig_to_pil(fig)
    plt.close(fig)
    return img


def img_scope(w, h) -> Image.Image:
    fig, ax = make_fig(w, h)

    # Summary bar
    summary = (f"Carve-Out Model: {P['model']}   |   "
               f"{P['employees']} employees   {P['sites']} sites   "
               f"{P['apps']} apps   {P['devices']} devices   Regions: {P['regions']}   "
               f"TSA: {P['tsa']} post Day 1")
    rect = mpatches.FancyBboxPatch(
        (0.0, 0.90), 1.0, 0.10, boxstyle="square,pad=0",
        linewidth=0, facecolor=BLUE, transform=ax.transAxes,
    )
    ax.add_patch(rect)
    ax.text(0.02, 0.958, summary, color=WHITE, fontsize=9, fontfamily=FONT,
            fontweight="bold", va="center", transform=ax.transAxes)

    # Two columns: In Scope / Out of Scope
    col1_x, col2_x = 0.01, 0.60

    draw_section_header(ax, 0.87, "IN SCOPE - IT Workstreams", x0=col1_x, fontsize=10)
    in_scope = [
        "IT Infrastructure  - WAN/LAN (180 sites), AD, servers, M365/Azure, telephony",
        "Commercial IT / ERP  - SAP migration (shell copy), CRM, FSM, BPO",
        "Other Applications  - ~500 apps: SharePoint, Confluence/Jira, eng. tools",
        "Engineering IT  - PLM/Windchill, FOSS compliance, developer network",
        "Production IT  - OT Security, MES, DOT, plant telephony",
        "HR IT  - HR systems for 8,000 employees (country-by-country rules)",
        "IT Org & Processes  - TOM, ITSM/ServiceNow, ITO contracting",
        "IT Contracts/Licences  - Change of control, SAM, ~6,000 device licences",
        "IT Security  - IAM/Saviynt, CISO, ISO 27001, BCM, GDPR",
    ]
    draw_bullet_block(ax, 0.82, in_scope, x0=col1_x + 0.01, fontsize=9,
                      line_gap=0.082, max_chars=62)

    draw_section_header(ax, 0.87, "OUT OF SCOPE", x0=col2_x, fontsize=10)
    out_scope = [
        "Business process redesign (IT separation only)",
        "Post-TSA operational optimisation (JCI IT post-closure)",
        "Non-IT business functions (separate functional workstreams)",
    ]
    draw_bullet_block(ax, 0.82, out_scope, x0=col2_x + 0.01, fontsize=9,
                      line_gap=0.082, max_chars=40)

    draw_section_header(ax, 0.58, "Country-Specific Complexity", x0=col2_x, fontsize=10)
    country = [
        "Brazil  - ERP tax law complexity (extreme)",
        "China   - Local FTS / customs systems",
        "India   - Local IdM policies (RBIN)",
        "Mexico  - Legal entity delay risk",
    ]
    draw_bullet_block(ax, 0.53, country, x0=col2_x + 0.01, fontsize=9,
                      line_gap=0.082, max_chars=40)

    # TSA note
    draw_section_header(ax, 0.15, "TSA Framework", x0=col2_x, fontsize=10)
    ax.text(col2_x + 0.01, 0.10,
            f"JCI operates IT for the carved-out heating business for {P['tsa']} post-Day 1 ({P['day1']}).\n"
            "TSA exit criteria to be agreed at QG1. Minimal scope preferred.",
            color=DKGRAY, fontsize=8.5, fontfamily=FONT, va="top",
            transform=ax.transAxes)

    plt.tight_layout(pad=0.1)
    img = fig_to_pil(fig)
    plt.close(fig)
    return img


def img_business_case(w, h) -> Image.Image:
    fig, ax = make_fig(w, h)

    # Strategic rationale
    draw_section_header(ax, 0.97, "Strategic Rationale", fontsize=10)
    rationale = [
        f"Bosch is acquiring JCI's heating business, requiring a full IT carve-out",
        f"SPA/APA obligations require full IT independence by Day 1 ({P['day1']})",
        f"Stand Alone carve-out: no shared platforms post-18-month TSA exit",
        f"KPMG methodology based on proven FRAME carve-out (reference engagement)",
    ]
    draw_bullet_block(ax, 0.90, rationale, fontsize=9.5, line_gap=0.07)

    # Financial table
    draw_section_header(ax, 0.63, f"Financial Summary  -  Labour Budget: {P['budget_labour']}", fontsize=10)
    fin_rows = [
        ["Category",          "Budget",          "% of Total", "Notes"],
        ["Governance / PMO",  P["budget_gov"],   "12.4%",      "Bosch IT PM + JCI IT PM"],
        ["Bosch IT Team",     P["budget_bosch"], "38.2%",      "Infra, ERP, Security, Apps WS leads"],
        ["JCI IT Team",       P["budget_jci"],   "23.5%",      "TSA operations, landscape assessment"],
        ["KPMG Consulting",   P["budget_kpmg"],  "25.5%",      "Methodology, WS leads, PM support"],
        ["TOTAL (Labour)",    P["budget_labour"],"100%",       "CFO sign-off required Apr 15, 2026"],
        ["Hardware/SW/Infra", "TBD",             "-",          "Separate CFO approval process"],
    ]
    draw_table(ax, fin_rows, [0.25, 0.17, 0.12, 0.45],
               y0=0.57, row_height=0.068, fontsize=9)

    # Value at risk
    draw_section_header(ax, 0.10, "Value at Risk  (if IT separation delays Day 1)", fontsize=10)
    risk_items = [
        "Legal / regulatory exposure from SPA breach",
        "Operational disruption across 180 sites, 8,000 employees",
        "TSA cost overrun: Bosch charges JCI for extended TSA services",
    ]
    draw_bullet_block(ax, 0.04, risk_items, fontsize=9, line_gap=0.055)

    plt.tight_layout(pad=0.1)
    img = fig_to_pil(fig)
    plt.close(fig)
    return img


def img_governance(w, h) -> Image.Image:
    fig, ax = make_fig(w, h)

    draw_section_header(ax, 0.97, "Governance Structure", fontsize=10)
    gov_rows = [
        ["Body",              "Frequency", "Members",                                    "Role"],
        ["Steering Committee","Monthly",   "Bosch CIO, JCI EVP, CFOs, Legal, Audit",    "QG approvals, budget, risk escalation"],
        ["PMO",               "Weekly",    "Bosch IT PM + JCI IT PM + KPMG Lead",        "Schedule, issue resolution, coordination"],
        ["WS Review",         "Bi-weekly", "9 Workstream Leads",                          "Deliverable tracking, blockers, dependencies"],
        ["Regional Stand-up", "Weekly",    "AP / AM / EMEA IT Leads",                    "Regional deployment progress"],
    ]
    draw_table(ax, gov_rows, [0.17, 0.12, 0.35, 0.37],
               y0=0.91, row_height=0.09, fontsize=9)

    draw_section_header(ax, 0.53, "Quality Gate Schedule", fontsize=10)
    qg_rows = [
        ["Gate",  "Date",      "Owner",              "Criteria"],
        ["QG 1",  P["qg1"],    "Steering Committee", "Governance approved; 180-site inventory; scope frozen"],
        ["QG 2",  P["qg2"],    "Steering Committee", "Architecture approved; migration strategy locked; wave plan done"],
        ["QG 3",  P["qg3"],    "Steering Committee", "Merger Zone built; SAP shell copy complete; WAN ordered"],
        ["QG 4",  P["qg4"],    "Steering Committee", "All migrations to Merger Zone complete; cutover ready"],
        ["Day 1", P["day1"],   "All",                "Big Bang cutover; SPA/APA signed; TSA activated"],
    ]
    draw_table(ax, qg_rows, [0.09, 0.13, 0.18, 0.60],
               y0=0.47, row_height=0.076, fontsize=9)

    draw_section_header(ax, 0.07, "Architecture Governance", fontsize=10)
    ax.text(0.02, 0.02,
            "Architecture Board: Bosch IT Architects + JCI IT Architects + KPMG  |  "
            "LeanIX as master CMDB  |  Frozen Zone from Signing - no production changes without IT PM approval",
            color=DKGRAY, fontsize=8.5, fontfamily=FONT, va="top", transform=ax.transAxes)

    plt.tight_layout(pad=0.1)
    img = fig_to_pil(fig)
    plt.close(fig)
    return img


def img_assumptions(w, h) -> Image.Image:
    fig, ax = make_fig(w, h)
    col1_x, col2_x = 0.01, 0.52

    draw_section_header(ax, 0.97, "Assumptions", x0=col1_x, fontsize=10)
    assumptions = [
        "A1  Carve-Out Model is Stand Alone (Steering to confirm May 2026)",
        "A2  18-month TSA is sufficient for full IT independence (FRAME benchmark)",
        "A3  Big Bang cutover on 01.11.2027 - no staggered go-live",
        "A4  KPMG is the primary consulting partner for all phases",
        "A5  JCI IT takes over TSA operations from Bosch IT at Day 1",
        "A6  LeanIX used as master application inventory (~500 apps)",
        "A7  Non-selective migration approach for data separation (mirrors FRAME)",
        "A8  3 regional DC hubs via co-locators (AP / AM / EMEA)",
    ]
    draw_bullet_block(ax, 0.90, assumptions, x0=col1_x + 0.01, fontsize=9,
                      line_gap=0.085, max_chars=52)

    draw_section_header(ax, 0.97, "Constraints & Dependencies", x0=col2_x, fontsize=10)
    constraints = [
        "C1  WAN ordering must start by Aug 2026 (4-6 month lead time - critical path)",
        "C2  SAP shell copy requires 9-12 months - must start in Concept phase",
        "C3  Co-locator / DC selection finalised by 13.01.2027",
        "C4  Frozen Zone from Signing - minimise changes to JCI production environment",
        "C5  Brazil ERP / China FTS / India RBIN - engage regional boards Month 2",
        "C6  IAM/Saviynt SC2 connector: restricted to 10 days before closing",
        "C7  SPA/APA legal timelines govern Day 1 date - IT must align",
        "C8  EUR 5.1M labour budget is fixed; hardware/SW needs separate approval",
        "D1  Legal entity ID needed before IT separation rules finalised",
        "D2  App inventory (LeanIX) complete before wave plan locked (QG2)",
    ]
    draw_bullet_block(ax, 0.90, constraints, x0=col2_x + 0.01, fontsize=9,
                      line_gap=0.085, max_chars=50)

    plt.tight_layout(pad=0.1)
    img = fig_to_pil(fig)
    plt.close(fig)
    return img


def img_dev_approach(w, h) -> Image.Image:
    fig, ax = make_fig(w, h)

    draw_section_header(ax, 0.97,
        "Delivery Approach: Classic / Phased (6-Phase Programme)  -  Big Bang Cutover",
        fontsize=10)
    ax.text(0.02, 0.89,
            "Rationale: Tightly coupled ERP and logistics processes across 180 sites make staggered "
            "go-live impractical. Big Bang mirrors the proven FRAME carve-out approach.",
            color=DKGRAY, fontsize=9, fontfamily=FONT, va="top", transform=ax.transAxes)

    phases = [
        ["Phase",    "Period",                  "Gate",      "Key Deliverables"],
        ["1  INIT",  f"{P['kickoff']} - {P['qg1']}", f"QG1 {P['qg1']}",
         "Governance, IT inventory (180 sites), project charter approved"],
        ["2  CONCEPT", f"{P['qg1']} - {P['qg2']}", f"QG2 {P['qg2']}",
         "IT architecture, migration strategy, wave plan, TSA service catalogue"],
        ["3  DEV",   f"{P['qg2']} - {P['qg3']}", f"QG3 {P['qg3']}",
         "3 regional DCs built, AD/M365, SAP shell copy, CMDB, WAN ordered"],
        ["4  IMPL",  f"{P['qg3']} - {P['qg4']}", f"QG4 {P['qg4']}",
         "Wave-based migrations complete: WAN, files, apps, client re-imaging"],
        ["5  GO-LIVE", f"{P['qg4']} - {P['day1']}", f"Day 1 {P['day1']}",
         "Big Bang cutover; independent operations; TSA activation; hypercare starts"],
        ["6  STAB",  f"{P['day1']} - {P['closure']}", f"Closure {P['closure']}",
         "3-month hypercare; TSA exit planning; lessons learned; project closure"],
    ]
    draw_table(ax, phases, [0.10, 0.23, 0.17, 0.50],
               y0=0.79, row_height=0.098, fontsize=9)

    draw_section_header(ax, 0.16, "Tooling & Methodology", fontsize=10)
    ax.text(0.02, 0.10,
            "Methodology: FRAME carve-out playbook (Bosch M&A IT) + KPMG M&A IT framework\n"
            "Tools: MS Project (schedule)   LeanIX (app inventory)   ServiceNow (ITSM)   SharePoint (PMO)   Saviynt (IAM)",
            color=DKGRAY, fontsize=9, fontfamily=FONT, va="top", transform=ax.transAxes)

    plt.tight_layout(pad=0.1)
    img = fig_to_pil(fig)
    plt.close(fig)
    return img


def img_risks(w, h) -> Image.Image:
    fig, ax = make_fig(w, h)

    draw_section_header(ax, 0.97,
        "Top 5 Risks  -  Rating = Probability x Impact (scale 1-25)", fontsize=10)

    risk_rows = [
        ["#", "Risk",                                               "P", "I", "Rating", "Owner",        "Mitigation"],
        ["R1", "Schedule Compression (WAN lead-time + SAP)",       "5", "5", "25",    "Bosch IT PM",
         "Order WAN by Month 4; SAP shell copy in parallel; agile UAT sprints"],
        ["R2", "Scope Creep (~500 apps; exception requests)",       "4", "5", "20",    "KPMG + WS Leads",
         "Lock inventory at QG1; wave-based prioritisation; Frozen Zone Month 12"],
        ["R3", "Data Separation Errors (non-selective migration)", "5", "5", "25",    "KPMG + Legal",
         "Finalise rules at QG1; third-party audit pre-cutover; spot-check 10% of records"],
        ["R4", "Country Complexity (Brazil / China / India)",       "4", "5", "20",    "Regional IT Leads",
         "Engage regional legal/finance Month 2; regional governance boards Month 3"],
        ["R5", "Resource Availability (Bosch/JCI competing)",      "3", "5", "15",    "IT PM",
         "Hybrid staffing; cross-training from Month 4; early SME contractor hiring"],
    ]

    # Draw table with colour-coded rating column
    cw = [0.04, 0.21, 0.04, 0.04, 0.08, 0.16, 0.44]
    x_positions = [0.01]
    for c in cw[:-1]:
        x_positions.append(x_positions[-1] + c)
    total_w = sum(cw)

    RATING_COLORS = {"25": RED, "20": AMBER, "15": "#DDAA00", "15": GREEN}

    for r_idx, row in enumerate(risk_rows):
        rh = 0.10 if r_idx > 0 else 0.065
        y_top = 0.91 - (0 if r_idx == 0 else 0.065 + (r_idx - 1) * 0.135)
        if r_idx > 0:
            y_top = 0.91 - 0.065 - (r_idx - 1) * 0.135

        y_bot = y_top - rh
        is_hdr = r_idx == 0
        bg = BLUE if is_hdr else (LTBLUE if r_idx % 2 == 0 else WHITE)
        fc = WHITE if is_hdr else DKGRAY
        fw = "bold" if is_hdr else "normal"

        rect = mpatches.FancyBboxPatch(
            (0.01, y_bot), total_w, rh,
            boxstyle="square,pad=0", linewidth=0.3,
            edgecolor="#AAAAAA", facecolor=bg,
            transform=ax.transAxes, clip_on=True,
        )
        ax.add_patch(rect)

        for c_idx, (cell, cx, cw_) in enumerate(zip(row, x_positions, cw)):
            color = fc
            fw2 = fw
            if not is_hdr and c_idx == 4:  # Rating col
                color = RED if cell == "25" else (AMBER if cell == "20" else "#997700")
                fw2 = "bold"
            ax.text(cx + 0.005, y_bot + rh * 0.5, str(cell),
                    color=color, fontsize=8 if c_idx == 6 else 8.5,
                    fontfamily=FONT, fontweight=fw2,
                    va="center", transform=ax.transAxes, clip_on=True)

    ax.text(0.01, 0.04,
            "Risk categories: ScR=Schedule  SR=Scope  RR=Resource  BtR=Budget  QR=Quality  BR=Business  LR=Legal  CR=Customer\n"
            "Full register: Trinity_Risk_Register.xlsx  |  Reviewed monthly at Steering Committee",
            color=DKGRAY, fontsize=8, fontfamily=FONT, va="top", transform=ax.transAxes)

    plt.tight_layout(pad=0.1)
    img = fig_to_pil(fig)
    plt.close(fig)
    return img


def img_timeline(w, h) -> Image.Image:
    """High-level milestone roadmap — 6 phases with quality gate diamonds."""
    from datetime import date as _date

    fig, ax = make_fig(w, h)

    START_YM  = (2026, 7)   # Jul 2026
    TOTAL_MO  = 20          # Jul 2026 → Feb 2028
    MARGIN_L  = 0.04
    MARGIN_R  = 0.03
    BAR_Y     = 0.52        # vertical centre of phase bars (axes fraction)
    BAR_H     = 0.19        # bar height

    def xf(yr, mo):
        """Convert year/month to x axes-fraction."""
        months = (yr - START_YM[0]) * 12 + mo - START_YM[1]
        return MARGIN_L + (months / TOTAL_MO) * (1 - MARGIN_L - MARGIN_R)

    # ── Phase bars ────────────────────────────────────────────────────────────
    phases = [
        ("Phase 1\nInitialization",  (2026, 7),  (2026, 10), "#004B87"),
        ("Phase 2\nConcept",         (2026, 10), (2027, 4),  "#0070C0"),
        ("Phase 3\nDevelopment",     (2027, 4),  (2027, 8),  "#0095C8"),
        ("Phase 4\nImplementation",  (2027, 8),  (2027, 11), "#00A3AD"),
        ("Phase 5\nGo-Live",         (2027, 11), (2027, 12), "#00703C"),
        ("Phase 6\nStabilization",   (2027, 12), (2028, 3),  "#404040"),
    ]
    for label, (sy, sm), (ey, em), color in phases:
        x_s = xf(sy, sm)
        x_e = xf(ey, em)
        rect = mpatches.FancyBboxPatch(
            (x_s, BAR_Y - BAR_H / 2), x_e - x_s, BAR_H,
            boxstyle="round,pad=0.004",
            facecolor=color, edgecolor=WHITE, linewidth=1.5,
            transform=ax.transAxes, zorder=2,
        )
        ax.add_patch(rect)
        cx = (x_s + x_e) / 2
        ax.text(cx, BAR_Y, label,
                color=WHITE, fontsize=7.5, fontfamily=FONT, fontweight="bold",
                ha="center", va="center", transform=ax.transAxes, zorder=3,
                multialignment="center")

    # ── Milestone diamonds ────────────────────────────────────────────────────
    # Alternate label placement above/below to avoid overlap for QG4/Day1
    milestones = [
        ("Kick-Off\n01.07.2026",  (2026, 7),  "#004B87", "above"),
        ("QG1\n30.09.2026",       (2026, 10), "#0070C0", "below"),
        ("QG2\n31.03.2027",       (2027, 4),  "#0095C8", "above"),
        ("QG3\n31.07.2027",       (2027, 8),  "#00A3AD", "below"),
        ("QG4\n31.10.2027",       (2027, 11), "#E08000", "above"),
        ("DAY 1\n01.11.2027",     (2027, 11), "#C00000", "below"),
        ("Closure\n28.02.2028",   (2028, 3),  "#404040", "above"),
    ]
    D_SIZE = 0.022
    for label, (yr, mo), color, pos in milestones:
        x = xf(yr, mo)
        bar_top = BAR_Y + BAR_H / 2
        bar_bot = BAR_Y - BAR_H / 2
        if pos == "above":
            dia_y   = bar_top + 0.09
            conn_y0 = bar_top
            lbl_y   = dia_y + D_SIZE + 0.03
            lbl_va  = "bottom"
        else:
            dia_y   = bar_bot - 0.09
            conn_y0 = bar_bot
            lbl_y   = dia_y - D_SIZE - 0.03
            lbl_va  = "top"
        # connector
        ax.plot([x, x], [conn_y0, dia_y], color=color, linewidth=1.0,
                transform=ax.transAxes, zorder=3)
        # diamond
        verts = [(x,          dia_y - D_SIZE),
                 (x + D_SIZE * 0.65, dia_y),
                 (x,          dia_y + D_SIZE),
                 (x - D_SIZE * 0.65, dia_y)]
        dia = plt.Polygon(verts, closed=True,
                          facecolor=color, edgecolor=WHITE, linewidth=0.8,
                          transform=ax.transAxes, zorder=4)
        ax.add_patch(dia)
        # label
        ax.text(x, lbl_y, label,
                color=color, fontsize=7, fontfamily=FONT, fontweight="bold",
                ha="center", va=lbl_va, transform=ax.transAxes, zorder=4,
                multialignment="center")

    # ── Time ruler ────────────────────────────────────────────────────────────
    ruler_y = BAR_Y - BAR_H / 2 - 0.30
    ax.plot([MARGIN_L, 1 - MARGIN_R], [ruler_y, ruler_y],
            color="#BBBBBB", linewidth=0.8, transform=ax.transAxes)
    for yr, mo in [(2026, 7), (2026, 10), (2027, 1), (2027, 4),
                   (2027, 7), (2027, 10), (2028, 1)]:
        x = xf(yr, mo)
        if x > 1 - MARGIN_R:
            break
        ax.plot([x, x], [ruler_y - 0.01, ruler_y], color="#BBBBBB",
                linewidth=0.8, transform=ax.transAxes)
        qtr = (mo - 1) // 3 + 1
        ax.text(x, ruler_y - 0.02, f"Q{qtr} {yr}",
                color="#777777", fontsize=6.5, fontfamily=FONT,
                ha="center", va="top", transform=ax.transAxes)

    # ── Title and footer ──────────────────────────────────────────────────────
    ax.text(0.5, 0.98,
            f"Project Trinity — High-Level Milestone Roadmap  ({P['kickoff']} – {P['closure']})",
            color=BLUE, fontsize=11, fontfamily=FONT, fontweight="bold",
            ha="center", va="top", transform=ax.transAxes)

    ax.text(0.5, 0.07,
            f"TSA Period: Day 1 ({P['day1']}) → Project Closure ({P['closure']}) — "
            f"JCI operates IT for carved-out heating business ({P['tsa']} TSA)",
            color=DKGRAY, fontsize=8, fontfamily=FONT,
            ha="center", va="top", transform=ax.transAxes)

    plt.tight_layout(pad=0.1)
    img = fig_to_pil(fig)
    plt.close(fig)
    return img


def img_cost(w, h) -> Image.Image:
    fig, ax = make_fig(w, h)

    draw_section_header(ax, 0.97, f"Cost / Budget Summary  -  Total Labour: {P['budget_labour']}", fontsize=10)

    cost_rows = [
        ["Category",            "Budget",           "% Labour", "Key Resources / Notes"],
        ["Governance / PMO",    P["budget_gov"],    "12.4%",
         "Bosch IT PM + JCI IT PM; SharePoint PMO; MS Project"],
        ["Bosch IT Team",       P["budget_bosch"],  "38.2%",
         "Infra, ERP/SAP, Security, Apps - 18 FTE across all phases"],
        ["JCI IT Team",         P["budget_jci"],    "23.5%",
         "TSA ops, landscape assessment - 12 FTE; ramp up Phases 3-4"],
        ["KPMG Consulting",     P["budget_kpmg"],   "25.5%",
         "Methodology, workstream leads, PM support - 10 FTE"],
        ["TOTAL (Labour)",      P["budget_labour"], "100%",
         "CFO sign-off required by Apr 15, 2026"],
        ["Hardware / SW / Infra", "TBD",            "-",
         "Separate CFO approval process; WAN, DC co-location, licences"],
        ["Travel & Vendor Fees",  "TBD",            "-",
         "Estimated separately per phase"],
    ]
    draw_table(ax, cost_rows, [0.21, 0.14, 0.12, 0.53],
               y0=0.89, row_height=0.082, fontsize=9)

    # Notes text (left side, below table)
    ax.text(0.02, 0.21,
            "Notes:\n"
            "  - Budget actuals tracked monthly vs cost plan\n"
            "  - KPMG engagement letter to be signed by 10.04.2026\n"
            "  - Hardware/SW budget TBD; subject to co-locator and WAN selection\n"
            "  - Budget report presented at every Steering Committee meeting",
            color=DKGRAY, fontsize=8.5, fontfamily=FONT, va="top",
            transform=ax.transAxes)

    # Mini bar chart (right side, below table — no overlap)
    cats  = ["Gov/PMO", "Bosch IT", "JCI IT", "KPMG"]
    vals  = [633.6, 1950, 1200, 1300]
    colors= [BLUE, "#005090", "#0090D0", "#00A0C0"]
    bar_ax = fig.add_axes([0.53, 0.03, 0.45, 0.20])
    bars = bar_ax.barh(cats, vals, color=colors, height=0.5)
    bar_ax.set_xlabel("Budget (kEUR)", fontsize=8, fontfamily=FONT)
    bar_ax.tick_params(labelsize=8)
    for bar, val in zip(bars, vals):
        bar_ax.text(bar.get_width() + 30, bar.get_y() + bar.get_height() / 2,
                    f"EUR {val/1000:.1f}M" if val >= 1000 else f"EUR {val:.0f}K",
                    va="center", fontsize=7.5, fontfamily=FONT)
    bar_ax.set_xlim(0, 2500)
    bar_ax.spines[["top", "right"]].set_visible(False)
    bar_ax.set_facecolor(WHITE)

    plt.tight_layout(pad=0.1)
    img = fig_to_pil(fig)
    plt.close(fig)
    return img


def img_resources(w, h) -> Image.Image:
    fig, ax = make_fig(w, h)

    draw_section_header(ax, 0.97,
        f"Resources  -  Peak Staffing: {P['ftes_peak']} during Phases 3-4", fontsize=10)

    res_rows = [
        ["Team",                "Phase 1", "Phase 2", "Phase 3-4", "Phase 5-6", "Budget"],
        ["Governance / PMO",    "2 FTE",   "2 FTE",   "2 FTE",    "2 FTE",    P["budget_gov"]],
        ["Bosch IT Team",       "8 FTE",   "14 FTE",  "18 FTE",   "10 FTE",   P["budget_bosch"]],
        ["JCI IT Team",         "4 FTE",   "8 FTE",   "12 FTE",   "6 FTE",    P["budget_jci"]],
        ["KPMG Consulting",     "4 FTE",   "8 FTE",   "10 FTE",   "4 FTE",    P["budget_kpmg"]],
        ["Regional IT Leads",   "1 FTE",   "2 FTE",   "3 FTE",    "2 FTE",    "(incl. above)"],
        ["TOTAL",               "19 FTE",  "34 FTE",  "45 FTE",   "24 FTE",   P["budget_labour"]],
    ]
    draw_table(ax, res_rows, [0.24, 0.12, 0.12, 0.13, 0.13, 0.17],
               y0=0.89, row_height=0.082, fontsize=9)

    draw_section_header(ax, 0.30, "Staffing Model & Confirmation", fontsize=10)
    model_items = [
        "Hybrid: Bosch IT (Phases 1-3 lead)  +  JCI IT (Phase 4-5 lead)  +  KPMG (all phases)",
        "Cross-training: Bosch to JCI knowledge transfer from Month 4",
        "External SME contractors: WAN, SAP migration, IAM/Saviynt - engaged from Month 3",
        "Resource ramp-up: Phases 1-2 (~19 FTE)  ->  Phases 3-4 (~45 FTE)  ->  Phase 6 (~24 FTE)",
        "Resources confirmed via resource plan and Steering Committee sign-off at QG1 (30.09.2026)",
        "Reference: Trinity_Resource_Summary.csv",
    ]
    draw_bullet_block(ax, 0.23, model_items, fontsize=9, line_gap=0.057)

    plt.tight_layout(pad=0.1)
    img = fig_to_pil(fig)
    plt.close(fig)
    return img


def img_status(w, h) -> Image.Image:
    fig, ax = make_fig(w, h)

    # Status banner
    rect = mpatches.FancyBboxPatch(
        (0.0, 0.88), 1.0, 0.12, boxstyle="square,pad=0",
        linewidth=0, facecolor=AMBER, transform=ax.transAxes,
    )
    ax.add_patch(rect)
    ax.text(0.02, 0.96, "PROJECT STATUS: PRE-LAUNCH - Awaiting Steering Committee Approvals",
            color=WHITE, fontsize=11, fontfamily=FONT, fontweight="bold",
            va="top", transform=ax.transAxes)
    ax.text(0.02, 0.90,
            f"Overall Health: ON TRACK   |   Report Date: {P['date']}   |   "
            f"Prepared by: Project Trinity Leadership Team",
            color=WHITE, fontsize=9, fontfamily=FONT, va="top", transform=ax.transAxes)

    phase_rows = [
        ["Phase",          "Gate Date",   "Status",          "Key Deliverables"],
        ["1 INITIALIZATION", P["qg1"],   "Ready to start",  "Governance model, IT inventory (180 sites), charter signed"],
        ["2 CONCEPT",        P["qg2"],   "Planned",         "IT architecture, migration strategy, app wave plan"],
        ["3 DEVELOPMENT",    P["qg3"],   "Planned",         "3 regional DCs built, Merger Zone online, SAP shell copy"],
        ["4 IMPLEMENTATION", P["qg4"],   "Planned",         "Wave-based migrations complete; WAN cutover done"],
        ["5 GO-LIVE",        P["day1"],  "CRITICAL PATH",   "Big Bang cutover; Day 1 independent ops; TSA activation"],
        ["6 STABILIZATION",  P["closure"],"Planned",        "Hypercare (3 months); TSA exit; project closure"],
    ]
    draw_table(ax, phase_rows, [0.19, 0.14, 0.14, 0.53],
               y0=0.82, row_height=0.082, fontsize=9)

    draw_section_header(ax, 0.28, "Decisions Required - Steering Committee Apr 15, 2026", fontsize=10)
    decisions = [
        f"1.  Approve Project Charter & Governance Model  (blocks team mobilisation - Jul 1 kick-off at risk)",
        f"2.  Authorise {P['budget_labour']} Labour Budget + Hardware/SW (TBD)  (resource hiring cannot proceed)",
        f"3.  Greenlight KPMG Engagement  (engagement letter by 10.04.2026)  (methodology support delayed)",
    ]
    draw_bullet_block(ax, 0.21, decisions, fontsize=9.5, line_gap=0.072)

    plt.tight_layout(pad=0.1)
    img = fig_to_pil(fig)
    plt.close(fig)
    return img


def img_comms(w, h) -> Image.Image:
    fig, ax = make_fig(w, h)

    draw_section_header(ax, 0.97, "Communication Plan", fontsize=10)

    comms_rows = [
        ["Audience",              "Frequency",  "Format",                     "Owner",          "Key Documents"],
        ["Steering Committee",    "Monthly",    "2-hr meeting + 1-page brief","IT PM + KPMG",   "Brief, Risk Register, OPL, Schedule"],
        ["Executive Leadership",  "Quarterly",  "3-5 page status report",     "IT PM",          "Executive Status Report, Dashboard"],
        ["IT Workstream Leads",   "Bi-weekly",  "WS review meeting",          "PMO",            "WS status, OPL items, dependencies"],
        ["All IT Staff",          "Monthly",    "Newsletter / SharePoint",    "PMO",            "Project update, upcoming milestones"],
        ["Legal / Finance",       "Ad hoc",     "Email + escalation call",    "IT PM + Legal",  "Contract status, risk alerts"],
        ["Regional Leads (AP/AM/EMEA)","Weekly","Stand-up call",              "Regional Leads", "Regional status, country blockers"],
        ["JCI Board",             "Quarterly",  "Executive dashboard",        "IT PM + CIO",    "Executive Dashboard HTML"],
    ]
    draw_table(ax, comms_rows, [0.18, 0.11, 0.18, 0.15, 0.37],
               y0=0.89, row_height=0.082, fontsize=8.5)

    draw_section_header(ax, 0.24, "Escalation Path", fontsize=10)
    ax.text(0.02, 0.18,
            "WS Lead  ->  IT PM  ->  Steering Committee  ->  Executive Leadership",
            color=BLUE, fontsize=10, fontfamily=FONT, fontweight="bold",
            va="top", transform=ax.transAxes)

    draw_section_header(ax, 0.12, "Confidentiality", fontsize=10)
    ax.text(0.02, 0.06,
            "All project documents: INTERNAL USE ONLY  |  "
            "Distribution limited to named stakeholders per this plan  |  "
            "External sharing requires Steering Committee approval",
            color=DKGRAY, fontsize=9, fontfamily=FONT, va="top",
            transform=ax.transAxes)

    plt.tight_layout(pad=0.1)
    img = fig_to_pil(fig)
    plt.close(fig)
    return img


def img_doc_history(w, h) -> Image.Image:
    fig, ax = make_fig(w, h)

    draw_section_header(ax, 0.97, "Document History", fontsize=10)
    hist_rows = [
        ["Version", "Date",       "Author",                          "Status / Changes"],
        ["1.0",     "29.03.2026", f"{P['pm']} / KPMG Lead",         "Initial draft - Pre-launch"],
        ["1.1",     "15.04.2026", "Steering Committee",              "Pending approval at SC meeting Apr 15"],
        ["",        "",           "",                                 ""],
    ]
    draw_table(ax, hist_rows, [0.10, 0.15, 0.30, 0.45],
               y0=0.89, row_height=0.10, fontsize=9.5)

    draw_section_header(ax, 0.55, "Template Information", fontsize=10)
    ax.text(0.02, 0.49,
            "Template:    Bosch IT Project Charter Template (FRAME M&A IT reference engagement)\n"
            "Generator:   generate_project_charter.py  |  Project Trinity  |  JCI to Bosch\n"
            "Methodology: FRAME carve-out playbook (Bosch M&A IT) + KPMG M&A IT framework",
            color=DKGRAY, fontsize=9.5, fontfamily=FONT, va="top", transform=ax.transAxes)

    plt.tight_layout(pad=0.1)
    img = fig_to_pil(fig)
    plt.close(fig)
    return img


# ── Slide-level title / breadcrumb updater ────────────────────────────────────

def _set_run_text(para, text: str) -> None:
    """Clear all runs on a paragraph and set a single new one with the given text."""
    from lxml import etree
    p = para._p
    for r in p.findall(qn("a:r")):
        p.remove(r)
    safe = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    r = etree.fromstring(
        '<a:r xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f"<a:t>{safe}</a:t></a:r>"
    )
    p.append(r)


def update_slide_titles(slide, new_title: str | None = None) -> None:
    """Update the title shape on a slide to new_title."""
    if not new_title:
        return
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.name in ("Titel 1", "Title 1", "Title 2"):
            tf = shape.text_frame
            if tf.paragraphs:
                _set_run_text(tf.paragraphs[0], new_title)
            break


# ── Slide content definitions ─────────────────────────────────────────────────
# Maps slide index (0-based) to (image_generator_fn, new_slide_title_or_None)

SLIDE_IMAGES = {
    1:  (img_agenda,        None),          # Slide 2: Agenda
    2:  (img_roles,         "Roles & Responsibilities"),
    3:  (img_goals,         "Goals & Success Criteria"),
    4:  (img_scope,         "Scope Description"),
    # slide 5 (index 5): Architecture - keep existing diagram, only label swap
    6:  (img_business_case, "Business Case"),
    7:  (img_governance,    "Governance"),
    8:  (img_assumptions,   "Assumptions & Constraints"),
    9:  (img_dev_approach,  "Development Approach"),
    10: (img_risks,         "Risks & Opportunities"),
    11: (img_timeline,      "Timeline & Milestones"),
    12: (img_cost,          "Cost / Budget"),
    13: (img_resources,     "Resources"),
    # slide 14 (index 14): Project Organisation - keep existing org chart
    15: (img_status,        "Status Report"),
    16: (img_comms,         "Communication Plan"),
    17: (img_doc_history,   "Document History"),
}

# ── Architecture / Timeline diagram label swaps ───────────────────────────────

ARCH_LABEL_MAP = {
    "UAES* uPortal":    "JCI IT Portal",
    "SUaaS":            "Bosch IT (TSA)",
    "Chatbot*":         "IAM / Saviynt",
    "SMT":              "ITSM / ServiceNow",
    "UIB*":             "Azure AD (JCI)",
    "MADAME":           "LeanIX / CMDB",
    "IDM":              "Bosch IdM (TSA exit)",
    "IT-Clearing PS/ORG": "JCI IT Operations",
    "UAES MDM":         "JCI MDM",
    "NC3":              "JCI Network",
    "UAES User":        "JCI Employee",
    "WF result":        "Access Result",
    "email notification": "Notification",
    "consumed system*": "JCI Application",
    "Link":             "TSA Link",
}

TIMELINE_LABEL_MAP = {
    "Project Pre.":                  "Pre-Launch",
    "Realization / Product Preparation": "Implementation",
    "Stabilization":                 "Stabilization",
    "Requirement Analysis":          "Landscape Assessment",
    "Cost Estimation":               "Architecture Design",
    "Project Plan":                  "Migration Strategy",
    "Resource Plan":                 "DC / WAN Build",
    "Architecture Check":            "Wave Migrations",
    "UAT":                           "Big Bang Cutover",
    "Feb. 2022":                     "Jul 2026",
    "Mar. 2022":                     "Oct 2026",
    "Apr. 2022":                     "Mar 2027",
    "May. 2022":                     "Jul 2027",
    "Jun. 2022":                     "Oct 2027",
    "July. 2022":                    "Nov 2027",
    "Aug. 2022":                     "Feb 2028",
    "Go-Live":                       "Day 1 Go-Live",
    "24.02.2022":                    P["kickoff"],
    "29.03.2022":                    P["date"],
    "01.08.2022":                    P["qg4"],
    "15.08.2022":                    P["day1"],
}


def swap_diagram_labels(slide, label_map: dict) -> None:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    for old, new in label_map.items():
                        if old in run.text:
                            run.text = run.text.replace(old, new)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Generate Project Trinity Project Charter PPTX")
    parser.add_argument("--out", default="Trinity_Project_Charter.pptx")
    args = parser.parse_args()

    output = BASE_DIR / args.out
    if not TEMPLATE.exists():
        raise FileNotFoundError(f"Template not found: {TEMPLATE}")

    print(f"Copying template -> {output.name}")
    shutil.copy2(TEMPLATE, output)

    prs = Presentation(output)

    print("Replacing project name in all text runs ...")
    replace_globally(prs)

    slides = list(prs.slides)
    print(f"Processing {len(slides)} slides ...")

    for idx, slide in enumerate(slides):
        print(f"  Slide {idx+1:2d}: ", end="", flush=True)

        if idx == 5:   # Architecture
            print("Architecture - swapping diagram labels")
            swap_diagram_labels(slide, ARCH_LABEL_MAP)
            update_slide_titles(slide, "Solution Architecture - Stand Alone Carve-Out Model")
            continue

        if idx == 14:  # Project Organisation
            print("Project Organisation - keeping org chart (label swap only)")
            org_map = {
                "Wu Chunliang(CI/MIR-AP) ": "Bosch IT PM",
                "Quality Manager":           "Quality & Risk Manager (KPMG)",
                " Yan Junhua (UAES/IT) (UD) |  HUANG Christina (CI/MIR-AP) (CI)":
                    " JCI IT PM  |  Bosch IT PM Lead",
            }
            swap_diagram_labels(slide, org_map)
            continue

        if idx not in SLIDE_IMAGES:
            print("(no image content - skipped)")
            continue

        img_fn, new_title = SLIDE_IMAGES[idx]
        print(f"{img_fn.__name__}()", end=" ", flush=True)

        # Get content area position and clear placeholder text
        left, top, width, height = hide_placeholder(slide, ph_idx=1)

        # Slide 4 (Goals / Target Cross): use the Group 33 exact dimensions
        if idx == 3:
            left, top, width, height = 0.535, 1.196, 11.095, 4.919

        # Agenda slide (idx=1) uses a 3-column layout; expand to full content width
        if idx == 1:
            left, width = 0.283, 11.43
            # also clear the extra non-standard text boxes on that slide
            for shape in slide.shapes:
                try:
                    ph = shape.placeholder_format
                    if ph is None and shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                run.text = ""
                except Exception:
                    pass

        # Generate the content image
        img = img_fn(width, height)

        # Insert image at placeholder position
        insert_image(slide, img, left, top, width, height)

        # Update slide title if provided
        if new_title:
            update_slide_titles(slide, new_title)

        print(f"-> image {int(width*DPI)}x{int(height*DPI)}px inserted")

    print(f"\nSaving -> {output}")
    prs.save(output)
    print(f"\nDone: {output}")
    print(f"  Project:  {P['name']} - {P['full']}")
    print(f"  Duration: {P['duration']}")
    print(f"  Budget:   {P['budget_labour']}")
    print(f"  Slides:   {len(slides)}")


if __name__ == "__main__":
    from datetime import datetime as _dt
    _t0 = _dt.now()
    print(f"Started : {_t0.strftime('%Y-%m-%d %H:%M:%S')}")
    try:
        main()
    finally:
        _t1 = _dt.now()
        print(f"Finished: {_t1.strftime('%Y-%m-%d %H:%M:%S')}  ({(_t1-_t0).total_seconds():.1f}s elapsed)")
