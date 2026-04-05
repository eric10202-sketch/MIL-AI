from __future__ import annotations

import base64
import csv
import os
import subprocess
import sys
import tempfile
from dataclasses import dataclass
from datetime import date, datetime
from pathlib import Path

from openpyxl import Workbook, load_workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas

HERE = Path(__file__).parent
OUT_DIR = HERE / "Charlie"
OUT_DIR.mkdir(exist_ok=True)

PROJECT = {
    "name": "Charlie",
    "seller": "Robert Bosch GmbH",
    "buyer": "Undisclosed (confidential)",
    "business": "Packaging",
    "model": "Stand Alone",
    "pm": "Gill Amandeep Singh (BD/MIL-PSM1)",
    "pmo": "KPMG",
    "sites": 37,
    "users": "3500+",
    "apps": 208,
    "sap_in_scope": True,
    "tsa_relevant": True,
    "start": "04/01/26",
    "golive": "06/01/27",
    "completion": "10/31/27",
}

SCHEDULE_XLSX = OUT_DIR / "Charlie_Project_Schedule.xlsx"
SCHEDULE_XML = OUT_DIR / "Charlie_Project_Schedule.xml"
RISK_XLSX = OUT_DIR / "Charlie_Risk_Register.xlsx"
COST_XLSX = OUT_DIR / "Charlie_Cost_Plan.xlsx"
CHARTER_HTML = OUT_DIR / "Charlie_Project_Charter.html"
EXEC_HTML = OUT_DIR / "Charlie_Executive_Dashboard.html"
KPI_HTML = OUT_DIR / "Charlie_Management_KPI_Dashboard.html"
PPTX_PATH = OUT_DIR / "Charlie_Management_Stakeholder_Presentation.pptx"

LOGO = HERE / "Bosch.png"
RISK_TEMPLATE = HERE / "BD_Risk-Register_template_en_V1.0_Dec2023.xlsx"
PPT_TEMPLATE = HERE / "Bosch presentation template.pptx"


@dataclass
class Task:
    id: int
    outline: int
    name: str
    duration: str
    start: str
    finish: str
    predecessors: str
    resources: str
    notes: str
    milestone: str


TASKS = [
    Task(1, 1, "Phase 0 - Initialization", "45 days", "04/01/26", "05/31/26", "", "PMO + Seller IT", "Mobilize carve-out governance and scope baseline", "No"),
    Task(2, 2, "0.1 Governance and Mobilization", "20 days", "04/01/26", "04/28/26", "", "PMO + Legal", "Project setup and governance", "No"),
    Task(3, 3, "Kickoff and RACI alignment", "8 days", "04/01/26", "04/10/26", "", "PMO + Workstream Leads", "Define ownership and cadence", "No"),
    Task(4, 3, "Scope baseline: SAP + 208 applications", "10 days", "04/11/26", "04/24/26", "3", "Enterprise Architecture", "Confirm in-scope systems", "No"),
    Task(5, 3, "TSA service domain framing", "5 days", "04/25/26", "04/30/26", "4", "Service Mgmt + Legal", "Define TSA-relevant service towers", "No"),
    Task(6, 2, "0.2 Discovery and dependency mapping", "23 days", "04/29/26", "05/31/26", "4", "SAP + App Owners", "Dependency map across SAP and non-SAP", "No"),
    Task(7, 3, "SAP landscape and interface inventory", "12 days", "04/29/26", "05/15/26", "4", "SAP CoE", "Instances, interfaces and batch jobs", "No"),
    Task(8, 3, "Application wave framing (208 apps)", "12 days", "05/16/26", "05/31/26", "7", "App Mgmt", "Wave and criticality matrix", "No"),
    Task(9, 2, "QG0 - Initialization Quality Gate", "0 days", "05/31/26", "05/31/26", "5,8", "Steering Committee", "Mandatory gate", "Yes"),
    Task(10, 1, "Phase 1 - Concept", "65 days", "06/01/26", "08/31/26", "9", "All Workstreams", "Target operating concept", "No"),
    Task(11, 2, "1.1 Target architecture concept", "35 days", "06/01/26", "07/19/26", "9", "Architecture + Security", "Stand-alone target landscape", "No"),
    Task(12, 3, "Standalone SAP concept", "15 days", "06/01/26", "06/21/26", "9", "SAP CoE", "Future state SAP design", "No"),
    Task(13, 3, "Identity and workplace concept", "10 days", "06/22/26", "07/05/26", "12", "Infra + Security", "User lifecycle and access", "No"),
    Task(14, 3, "Data ownership and compliance concept", "10 days", "07/06/26", "07/19/26", "13", "Data Gov + Legal", "Retention and transfer rules", "No"),
    Task(15, 2, "1.2 TSA operating concept", "25 days", "07/20/26", "08/23/26", "14", "Service Mgmt + Buyer Liaison", "Service catalogue and exit approach", "No"),
    Task(16, 3, "Define TSA SLAs and support scope", "12 days", "07/20/26", "08/05/26", "14", "Service Mgmt", "Critical service levels", "No"),
    Task(17, 3, "Buyer readiness checkpoints", "10 days", "08/06/26", "08/19/26", "16", "PMO + Buyer Liaison", "Readiness criteria", "No"),
    Task(18, 3, "TSA commercial and governance baseline", "3 days", "08/20/26", "08/23/26", "17", "Finance + Legal", "Commercial principles", "No"),
    Task(19, 2, "QG1 - Concept Quality Gate", "0 days", "08/31/26", "08/31/26", "15,18", "Steering Committee", "Mandatory concept approval", "Yes"),
    Task(20, 1, "Phase 2 - Architecture and Design", "85 days", "09/01/26", "11/30/26", "19", "All Workstreams", "Detailed design and cutover planning", "No"),
    Task(21, 2, "2.1 SAP and integration design", "45 days", "09/01/26", "10/31/26", "19", "SAP + Integration", "Detailed build design", "No"),
    Task(22, 3, "SAP module and data migration design", "20 days", "09/01/26", "09/28/26", "19", "SAP + Data", "Object-level migration design", "No"),
    Task(23, 3, "Interface and API design", "15 days", "09/29/26", "10/19/26", "22", "Integration", "SAP and app integration contracts", "No"),
    Task(24, 3, "Test strategy and cutover blueprint", "10 days", "10/20/26", "10/31/26", "23", "Test Mgmt + PMO", "Integration/UAT strategy", "No"),
    Task(25, 2, "2.2 Non-SAP application design (208)", "45 days", "09/15/26", "11/15/26", "19", "App Owners", "Wave design and decommission paths", "No"),
    Task(26, 3, "Wave 1 and 2 design", "20 days", "09/15/26", "10/12/26", "19", "App Mgmt", "Business-critical apps", "No"),
    Task(27, 3, "Wave 3 and 4 design", "20 days", "10/13/26", "11/09/26", "26", "App Mgmt", "Long-tail applications", "No"),
    Task(28, 3, "Deployment and rollback planning", "5 days", "11/10/26", "11/15/26", "27", "PMO + Ops", "Cutover fallback options", "No"),
    Task(29, 2, "2.3 Infrastructure and security design", "40 days", "10/01/26", "11/30/26", "19", "Infra + Security", "Global site architecture", "No"),
    Task(30, 2, "QG2/3 - Combined Design & Build Gate", "0 days", "11/30/26", "11/30/26", "24,28,29", "Steering Committee", "Combined mandatory gate", "Yes"),
    Task(31, 1, "Phase 3 - Development, Build and Test", "110 days", "12/01/26", "04/30/27", "30", "All Workstreams", "Build, migration and testing", "No"),
    Task(32, 2, "3.1 SAP build and data migration", "70 days", "12/01/26", "02/28/27", "30", "SAP CoE", "Build and validate SAP", "No"),
    Task(33, 2, "3.2 Application migration waves", "85 days", "12/15/26", "03/31/27", "30", "App Teams", "Migrate 208 applications", "No"),
    Task(34, 2, "3.3 Infra build and workplace rollout", "75 days", "12/01/26", "03/15/27", "30", "Infra + Workplace", "37-site readiness", "No"),
    Task(35, 2, "3.4 Integration and UAT", "55 days", "02/15/27", "04/30/27", "32,33,34", "Test Team + Business", "End-to-end readiness", "No"),
    Task(36, 2, "Operational readiness review", "10 days", "04/17/27", "04/30/27", "35", "PMO + Ops", "Go/no-go review", "No"),
    Task(37, 1, "Phase 4 - GoLive and Closure", "130 days", "05/01/27", "10/31/27", "36", "All Workstreams", "GoLive, TSA stabilization and closure", "No"),
    Task(38, 2, "4.1 Pre-GoLive finalization", "20 days", "05/01/27", "05/31/27", "36", "PMO + Ops", "Final readiness", "No"),
    Task(39, 2, "QG4 - Pre-GoLive Quality Gate", "0 days", "05/31/27", "05/31/27", "38", "Steering Committee", "Mandatory pre-GoLive gate", "Yes"),
    Task(40, 2, "GoLive Day 1", "0 days", "06/01/27", "06/01/27", "39", "All Workstreams", "GoLive after QG4", "Yes"),
    Task(41, 2, "4.2 Hypercare and TSA-supported stabilization", "65 days", "06/02/27", "08/31/27", "40", "Ops + Service Mgmt", "Stabilize and transition", "No"),
    Task(42, 2, "4.3 TSA exit and buyer handover", "22 days", "09/01/27", "09/30/27", "41", "Buyer IT + PMO", "Service handover", "No"),
    Task(43, 2, "4.4 Financial close and archive", "23 days", "10/01/27", "10/31/27", "42", "PMO + Finance", "Closure pack", "No"),
    Task(44, 2, "QG5 - Programme Closure Quality Gate", "0 days", "10/31/27", "10/31/27", "43", "Steering Committee", "Programme closure", "Yes"),
]


def write_schedule():
    headers = [
        "ID", "Outline Level", "Name", "Duration", "Start", "Finish",
        "Predecessors", "Resource Names", "Notes", "Milestone",
    ]
    wb = Workbook()
    ws = wb.active
    ws.title = "Project Schedule"

    fills = {
        "header": PatternFill("solid", fgColor="002147"),
        "phase": PatternFill("solid", fgColor="003B6E"),
        "section": PatternFill("solid", fgColor="0066CC"),
        "milestone": PatternFill("solid", fgColor="FFF2CC"),
        "alt": PatternFill("solid", fgColor="EFF4FB"),
    }
    thin = Side(style="thin", color="CCCCCC")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    widths = [6, 8, 52, 12, 12, 12, 14, 28, 48, 10]
    for i, w in enumerate(widths, start=1):
        ws.column_dimensions[get_column_letter(i)].width = w

    for c, h in enumerate(headers, start=1):
        cell = ws.cell(1, c, h)
        cell.font = Font(name="Calibri", size=10, bold=True, color="FFFFFF")
        cell.fill = fills["header"]
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell.border = border

    for r, t in enumerate(TASKS, start=2):
        vals = [t.id, t.outline, t.name, t.duration, t.start, t.finish, t.predecessors, t.resources, t.notes, t.milestone]
        is_phase = t.outline == 1
        is_section = t.outline == 2
        is_milestone = t.milestone == "Yes"

        for c, v in enumerate(vals, start=1):
            cell = ws.cell(r, c)
            if c == 3:
                indent = "" if t.outline == 1 else ("  " if t.outline == 2 else "    ")
                cell.value = indent + str(v)
            else:
                cell.value = v
            cell.border = border
            cell.alignment = Alignment(vertical="center", wrap_text=(c in (3, 8, 9)))
            cell.font = Font(name="Calibri", size=9, bold=is_phase or is_section or is_milestone, color="FFFFFF" if is_phase or (is_section and not is_milestone) else "000000")

        if is_phase:
            fill = fills["phase"]
        elif is_milestone:
            fill = fills["milestone"]
        elif is_section:
            fill = fills["section"]
        elif r % 2 == 0:
            fill = fills["alt"]
        else:
            fill = None
        if fill:
            for c in range(1, len(headers) + 1):
                ws.cell(r, c).fill = fill

    ws.freeze_panes = "A2"
    ws.auto_filter.ref = f"A1:{get_column_letter(len(headers))}1"
    wb.save(SCHEDULE_XLSX)

    fd, tmp_csv = tempfile.mkstemp(suffix=".csv")
    os.close(fd)
    try:
        with open(tmp_csv, "w", newline="", encoding="utf-8") as f:
            w = csv.writer(f)
            w.writerow(headers)
            for t in TASKS:
                w.writerow([t.id, t.outline, t.name, t.duration, t.start, t.finish, t.predecessors, t.resources, t.notes, t.milestone])

        result = subprocess.run(
            [sys.executable, str(HERE / "generate_msp_xml.py"), "--csv", tmp_csv, "--out", str(SCHEDULE_XML), "--project", "Project Charlie"],
            capture_output=True,
            text=True,
        )
        if result.returncode != 0:
            raise RuntimeError(result.stderr)
    finally:
        try:
            os.unlink(tmp_csv)
        except OSError:
            pass


def write_risk_register():
    wb = load_workbook(RISK_TEMPLATE)
    info = wb["Info"]
    rr = wb["Risk Register"]

    info["C4"] = "Charlie-RR-001"
    info["C5"] = "Charlie_Risk_Register.xlsx"
    info["C6"] = "Project Charlie - Packaging IT Carve-Out"
    info["C7"] = "Charlie"
    info["C8"] = PROJECT["pm"]

    risks = [
        ("Schedule", "Delivery timeline compression before GoLive", "GoLive delay", "Complex dependencies across SAP and 208 apps", "PMO", "Internal", "High", "70%", "Mitigate", "Wave-level planning, weekly control tower", "05/31/27", "in progress"),
        ("Engineering", "Undiscovered dependencies in non-SAP app estate", "Migration rework", "Legacy integrations not documented", "App Mgmt", "Internal", "High", "50%", "Mitigate", "Discovery sprints and integration maps", "10/31/26", "in progress"),
        ("Resources", "Specialist SAP skills shortage during build", "Slower build and testing", "Competing programs", "SAP Lead", "Internal", "Moderate", "50%", "Mitigate", "Early staffing lock and external backup", "12/15/26", "not started"),
        ("Legal & Compliance", "Data transfer approvals delayed", "Cutover delay", "Cross-border constraints", "Legal", "External", "High", "30%", "Mitigate", "Early legal review and approvals", "02/15/27", "in progress"),
        ("Budget", "Extended TSA duration cost growth", "Budget overrun", "Buyer readiness slower than expected", "Finance", "Internal", "Moderate", "50%", "Mitigate", "Exit criteria and monthly readiness checks", "09/30/27", "in progress"),
        ("Quality", "UAT defects exceed threshold", "Readiness at risk", "Complex test matrix", "Test Lead", "Internal", "High", "50%", "Mitigate", "Defect triage and hard freeze windows", "04/20/27", "in progress"),
        ("Strategy & Portfolio", "Late scope changes from business", "Plan churn", "Unfrozen requirements", "Business Lead", "Internal", "Moderate", "30%", "Avoid", "Scope freeze at QG1", "08/31/26", "not started"),
        ("Security", "IAM rollout delays for 3500+ users", "Day-1 access incidents", "Provisioning complexity", "Security Lead", "Internal", "High", "50%", "Mitigate", "Pilot provisioning and fallback runbooks", "05/15/27", "in progress"),
        ("Schedule", "Infrastructure readiness lag at remote sites", "Partial regional delay", "Carrier provisioning lead times", "Infra Lead", "External", "Moderate", "30%", "Transfer", "Carrier escalation and alternative links", "03/31/27", "in progress"),
        ("Engineering", "SAP cutover rehearsal incomplete", "Production instability", "Limited windows", "Cutover Manager", "Internal", "Very High", "30%", "Mitigate", "Two full rehearsals before QG4", "05/20/27", "in progress"),
    ]

    start_row = 5
    for idx, r in enumerate(risks, start=1):
        row = start_row + idx - 1
        category, event, effects, causes, owner, source, impact, prob, strategy, measure, due, status = r
        rr.cell(row, 2, f"CR-{idx:03d}")
        rr.cell(row, 3, datetime.today().strftime("%d-%b-%y"))
        rr.cell(row, 4, category)
        rr.cell(row, 5, causes)
        rr.cell(row, 6, event)
        rr.cell(row, 7, effects)
        rr.cell(row, 8, PROJECT["golive"])
        rr.cell(row, 9, owner)
        rr.cell(row, 10, source)
        rr.cell(row, 12, impact)
        rr.cell(row, 13, f'=_xlfn.IFNA(VLOOKUP(L{row},$D$182:$E$186,2,FALSE),"")')
        rr.cell(row, 14, prob)
        rr.cell(row, 15, f'=_xlfn.IFNA(VLOOKUP(N{row},$D$189:$E$193,2,FALSE),"")')
        rr.cell(row, 16, "Threat")
        rr.cell(row, 17, f"=M{row}*O{row}")
        rr.cell(row, 18, effects)
        rr.cell(row, 22, strategy)
        rr.cell(row, 23, measure)
        rr.cell(row, 24, due)
        rr.cell(row, 26, status)
        rr.cell(row, 27, datetime.today().strftime("%d-%b-%y"))
        rr.cell(row, 28, impact)
        rr.cell(row, 29, f'=_xlfn.IFNA(VLOOKUP(AB{row},$D$182:$E$186,2,FALSE),"")')
        rr.cell(row, 30, prob)
        rr.cell(row, 31, f'=_xlfn.IFNA(VLOOKUP(AD{row},$D$189:$E$193,2,FALSE),"")')
        rr.cell(row, 32, f"=AC{row}")
        rr.cell(row, 33, f"=AE{row}")
        rr.cell(row, 34, f"=AF{row}*AG{row}")
        rr.cell(row, 35, "Generated for Project Charlie")

    wb.save(RISK_XLSX)


def write_cost_plan():
    wb = Workbook()
    ws = wb.active
    ws.title = "Cost Plan"

    rows = [
        ["PROJECT CHARLIE - Cost Plan (Labour Only)", "", "", "", "", ""],
        ["Category / Task", "Resource", "Days", "Hours", "Rate EUR/h", "Cost EUR"],
        ["Metadata", "Seller: Robert Bosch GmbH | Buyer: Undisclosed | PM: Gill Amandeep Singh | PMO: KPMG", "", "", "", ""],
        ["Metadata", "Scope: SAP + 208 applications | 37 sites | 3500+ users | TSA relevant", "", "", "", ""],
        ["Metadata", "Timeline: 01 Apr 2026 to 31 Oct 2027", "", "", "", ""],
        ["Programme Management", "PMO and governance", 150, 1200, 125, 150000],
        ["SAP Separation", "SAP architecture/build/migration", 220, 1760, 135, 237600],
        ["Application Migration", "208 app migration waves", 260, 2080, 110, 228800],
        ["Infrastructure and Security", "37-site infrastructure and IAM", 180, 1440, 105, 151200],
        ["Testing and Cutover", "UAT, rehearsal, cutover", 120, 960, 100, 96000],
        ["Hypercare and TSA Transition", "Stabilization and handover", 140, 1120, 100, 112000],
        ["TOTAL", "", "", "", "", 975600],
        ["CAPEX", "Licensing, network and hardware", "", "", "", "TBC - to be approved at QG1"],
    ]

    dark = PatternFill("solid", fgColor="003B6E")
    mid = PatternFill("solid", fgColor="0066CC")
    light = PatternFill("solid", fgColor="EFF4FB")
    thin = Side(style="thin", color="CCCCCC")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    for col, w in zip("ABCDEF", [45, 55, 9, 9, 12, 14]):
        ws.column_dimensions[col].width = w

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row, start=1):
            cell = ws.cell(r, c, val)
            cell.border = border
            if r == 1:
                cell.fill = dark
                cell.font = Font(bold=True, color="FFFFFF", size=12)
            elif r == 2:
                cell.fill = mid
                cell.font = Font(bold=True, color="FFFFFF")
            elif 3 <= r <= 5:
                cell.fill = light
                cell.font = Font(bold=False)
            elif r == 12:
                cell.fill = PatternFill("solid", fgColor="C6D4E8")
                cell.font = Font(bold=True)
            elif r == 13:
                cell.fill = PatternFill("solid", fgColor="FFF2CC")
                cell.font = Font(bold=True)
        ws.row_dimensions[r].height = 18

    ws.merge_cells("A1:F1")
    ws.merge_cells("B3:F3")
    ws.merge_cells("B4:F4")
    ws.merge_cells("B5:F5")
    ws.freeze_panes = "A3"
    ws.auto_filter.ref = "A2:F2"
    wb.save(COST_XLSX)


def write_charter():
    logo = base64.b64encode(LOGO.read_bytes()).decode("ascii") if LOGO.exists() else ""
    html = f"""<!doctype html>
<html><head><meta charset='utf-8'><title>Charlie Project Charter</title>
<style>
body{{font-family:Calibri,Arial,sans-serif;background:#f0f3f8;color:#1a1a1a;margin:0}}
.container{{max-width:1000px;margin:0 auto;padding:24px}}
.hero{{background:linear-gradient(140deg,#002147,#003b6e,#0066cc);color:#fff;padding:28px;border-radius:14px}}
.bosch-logo{{display:flex;align-items:center;background:#fff;padding:4px 8px;border-radius:4px;width:fit-content;margin-bottom:10px}}
.card{{background:#fff;padding:18px;border-radius:10px;margin-top:14px}}
h2{{color:#003b6e;margin:0 0 10px 0}}
</style></head>
<body><div class='container'>
<div class='hero'>
<div class='bosch-logo'><img src='data:image/png;base64,{logo}' style='height:36px'></div>
<h1>Project Charlie - Project Charter</h1>
<p>Packaging IT carve-out from Robert Bosch GmbH to an undisclosed buyer</p>
</div>
<div class='card'><h2>Project Facts</h2>
<p><b>Seller:</b> {PROJECT['seller']}<br><b>Buyer:</b> {PROJECT['buyer']}<br><b>Business:</b> {PROJECT['business']}<br><b>Model:</b> {PROJECT['model']}<br><b>PM:</b> {PROJECT['pm']}<br><b>PMO:</b> {PROJECT['pmo']}</p>
<p><b>Scope:</b> SAP in scope, {PROJECT['apps']} applications, {PROJECT['sites']} worldwide sites, {PROJECT['users']} users, TSA relevant.</p>
<p><b>Timeline:</b> Start 01 Apr 2026, GoLive 01 Jun 2027, Completion 31 Oct 2027.</p>
</div>
<div class='card'><h2>Quality Gates</h2>
<p>QG0 (31 May 2026), QG1 (31 Aug 2026), QG2/3 combined (30 Nov 2026), QG4 Pre-GoLive (31 May 2027), QG5 Programme Closure (31 Oct 2027).</p>
</div>
</div></body></html>"""
    CHARTER_HTML.write_text(html, encoding="utf-8")


def write_exec_dashboard():
    today = date.today()
    golive = date(2027, 6, 1)
    qg4 = date(2027, 5, 31)
    html = f"""<!doctype html><html><head><meta charset='utf-8'><title>Charlie Executive Dashboard</title>
<style>body{{font-family:Calibri,Arial,sans-serif;background:#f4f6f9;margin:0}}.wrap{{max-width:1100px;margin:auto;padding:20px}}.hero{{background:linear-gradient(130deg,#002147,#003b6e,#0066cc);color:#fff;padding:20px;border-radius:12px}}.grid{{display:grid;grid-template-columns:repeat(4,1fr);gap:10px;margin-top:12px}}.k{{background:#fff;border-radius:8px;padding:12px}}</style></head>
<body><div class='wrap'><div class='hero'><h1>Project Charlie - Executive Dashboard</h1><p>Packaging carve-out | Stand Alone | TSA relevant | Dashboard date: {today.strftime('%d %b %Y')}</p></div>
<div class='grid'>
<div class='k'><b>GoLive</b><br>01 Jun 2027<br>{(golive-today).days} days</div>
<div class='k'><b>QG4 Pre-GoLive</b><br>31 May 2027<br>{(qg4-today).days} days</div>
<div class='k'><b>Scope</b><br>SAP + 208 apps</div>
<div class='k'><b>Footprint</b><br>37 sites | 3500+ users</div>
</div>
<div class='k' style='margin-top:12px'><b>Top Focus Areas:</b> SAP migration readiness, application wave execution, TSA service definition and exit, cutover rehearsal completion.</div>
</div></body></html>"""
    EXEC_HTML.write_text(html, encoding="utf-8")


def write_kpi_dashboard():
    today = date.today()
    start = date(2026, 4, 1)
    end = date(2027, 10, 31)
    elapsed = max(0, (today - start).days)
    total = (end - start).days
    pct = min(100, round((elapsed / total) * 100))
    html = f"""<!doctype html><html><head><meta charset='utf-8'><title>Charlie KPI Dashboard</title>
<style>body{{font-family:Calibri,Arial,sans-serif;background:#f4f6f9;margin:0}}.wrap{{max-width:1100px;margin:auto;padding:20px}}.hero{{background:#003b6e;color:#fff;padding:16px;border-radius:10px}}.grid{{display:grid;grid-template-columns:repeat(6,1fr);gap:10px;margin-top:12px}}.card{{background:#fff;padding:12px;border-radius:8px;text-align:center}}</style></head>
<body><div class='wrap'><div class='hero'><h2>Project Charlie - Management KPI Dashboard</h2><p>PM: {PROJECT['pm']} | PMO: {PROJECT['pmo']}</p></div>
<div class='grid'>
<div class='card'><b>SPI</b><br>1.00</div>
<div class='card'><b>CPI</b><br>1.00</div>
<div class='card'><b>Readiness</b><br>{pct}%</div>
<div class='card'><b>Apps Planned</b><br>{PROJECT['apps']}</div>
<div class='card'><b>Open Risks</b><br>10</div>
<div class='card'><b>Critical Issues</b><br>0</div>
</div>
<div class='card' style='margin-top:12px;text-align:left'><b>90-Day Priorities:</b> finalize TSA SLAs, complete detailed build design, lock infra readiness across 37 sites, execute migration wave rehearsals.</div>
</div></body></html>"""
    KPI_HTML.write_text(html, encoding="utf-8")


def write_monthly_pdf():
    report_name = f"Charlie_Monthly_Status_Report_{date.today().strftime('%b_%Y')}.pdf"
    out = OUT_DIR / report_name
    c = canvas.Canvas(str(out), pagesize=A4)
    w, h = A4
    c.setFillColor(colors.HexColor("#003b6e"))
    c.rect(0, h - 78, w, 78, fill=1, stroke=0)
    if LOGO.exists():
        c.drawImage(ImageReader(str(LOGO)), 24, h - 64, height=30, preserveAspectRatio=True, mask="auto")
    c.setFillColor(colors.white)
    c.setFont("Helvetica-Bold", 14)
    c.drawString(130, h - 38, "Project Charlie - Monthly Status Report")
    c.setFont("Helvetica", 9)
    c.drawString(130, h - 54, date.today().strftime("%d %b %Y"))

    y = h - 105
    c.setFillColor(colors.black)
    c.setFont("Helvetica", 10)
    lines = [
        "Seller: Robert Bosch GmbH | Buyer: Undisclosed (confidential)",
        "Business: Packaging | Model: Stand Alone | TSA: Relevant",
        "Scope: SAP in scope, 208 applications, 37 sites, 3500+ users",
        "Timeline: Start 01 Apr 2026 | GoLive 01 Jun 2027 | Completion 31 Oct 2027",
        "Milestones: QG0 31 May 2026, QG1 31 Aug 2026, QG2/3 30 Nov 2026, QG4 31 May 2027, QG5 31 Oct 2027",
        "Status summary: On baseline plan; major watchpoints are SAP migration readiness and TSA exit criteria.",
    ]
    for line in lines:
        c.drawString(30, y, line)
        y -= 18
    c.showPage()
    c.save()


def write_presentation():
    prs = Presentation(str(PPT_TEMPLATE)) if PPT_TEMPLATE.exists() else Presentation()
    while len(prs.slides) > 0:
        rid = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rid)
        del prs.slides._sldIdLst[0]

    cover = prs.slides.add_slide(prs.slide_layouts[0])
    if cover.shapes.title:
        cover.shapes.title.text = "Project Charlie"
    subtitle_text = "Packaging carve-out | Stand Alone | Report date: " + date.today().strftime("%d %b %Y")
    if len(cover.placeholders) > 1:
        cover.placeholders[1].text = subtitle_text

    def add_bullet_slide(title: str, subtitle: str, bullets: list[str]):
        s = prs.slides.add_slide(prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0])
        if s.shapes.title:
            s.shapes.title.text = title
        body = None
        for shp in s.shapes:
            if hasattr(shp, "text_frame") and shp != s.shapes.title:
                body = shp
                break
        if body is None:
            body = s.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.0), Inches(4.5))
        tf = body.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.bold = True
        p.font.size = Pt(18)
        for b in bullets:
            pp = tf.add_paragraph()
            pp.text = b
            pp.level = 0
            pp.font.size = Pt(14)

    add_bullet_slide("Executive Summary", "Project baseline and current outlook", [
        "Seller: Robert Bosch GmbH; Buyer: Undisclosed (confidential)",
        "Scope includes SAP and 208 applications",
        "Global footprint: 37 sites and 3500+ users",
        "TSA is relevant for transition and stabilization",
    ])
    add_bullet_slide("Timeline and Gates", "Planned quality-gate sequence", [
        "QG0: 31 May 2026",
        "QG1: 31 Aug 2026",
        "QG2/3 combined: 30 Nov 2026",
        "QG4 Pre-GoLive: 31 May 2027",
        "GoLive Day 1: 01 Jun 2027",
        "QG5 Programme Closure: 31 Oct 2027",
    ])
    add_bullet_slide("Risks and Decisions", "Management focus", [
        "Top risk themes: schedule compression, hidden app dependencies, TSA cost extension",
        "Decision needs: TSA governance baseline, resource protection, cutover rehearsal criteria",
    ])
    prs.save(str(PPTX_PATH))


def main():
    write_schedule()
    write_risk_register()
    write_cost_plan()
    write_charter()
    write_exec_dashboard()
    write_kpi_dashboard()
    write_monthly_pdf()
    write_presentation()
    print("Generated Charlie deliverables in:", OUT_DIR)


if __name__ == "__main__":
    main()
