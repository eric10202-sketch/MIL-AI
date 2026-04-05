#!/usr/bin/env python3
"""
Generate Trinity-CAM Management Stakeholder Presentation (PPTX).
Uses Bosch presentation template as base.
All content derived fresh from Trinity-CAM schedule/cost/risk data.
Structure reference: generate_bravo_stakeholder_presentation.py (FORMAT ONLY).
"""

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---- Bosch palette -----
BOSCH_NAVY  = RGBColor(0x00, 0x3B, 0x6E)
BOSCH_MID   = RGBColor(0x00, 0x51, 0x99)
BOSCH_BLUE  = RGBColor(0x00, 0x66, 0xCC)
BOSCH_LIGHT = RGBColor(0xE4, 0xED, 0xF9)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK   = RGBColor(0x1A, 0x1A, 0x1A)
MUTED       = RGBColor(0x5A, 0x64, 0x78)
GREEN       = RGBColor(0x00, 0x7A, 0x33)
AMBER       = RGBColor(0xE8, 0xA0, 0x00)
RED         = RGBColor(0xCC, 0x00, 0x00)

ROOT     = Path(__file__).resolve().parent
TEMPLATE = ROOT / "Reference" / "Bosch presentation template.pptx"
OUTPUT   = ROOT / "active-projects" / "Trinity-CAM" / "Trinity-CAM_Stakeholder_Presentation.pptx"
OUTPUT.parent.mkdir(parents=True, exist_ok=True)

REPORT_DATE = date.today().strftime("%d %b %Y")


# ─────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────

def create_prs() -> Presentation:
    if TEMPLATE.exists():
        return Presentation(str(TEMPLATE))
    print(f"Warning: template not found at {TEMPLATE}. Using blank presentation.")
    return Presentation()


def clear_slides(prs: Presentation) -> None:
    while len(prs.slides):
        r_id = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]


def set_ph(slide, idx, text) -> bool:
    for sh in slide.placeholders:
        if sh.placeholder_format.idx == idx:
            sh.text = text
            return True
    return False


def add_accent_bar(slide):
    bar = slide.shapes.add_shape(1, Inches(0.22), Inches(1.33), Inches(11.55), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BOSCH_BLUE
    bar.line.fill.background()


def add_corner_dot(slide):
    dot = slide.shapes.add_shape(9, Inches(11.65), Inches(0.25), Inches(0.22), Inches(0.22))
    dot.fill.solid()
    dot.fill.fore_color.rgb = BOSCH_BLUE
    dot.line.fill.background()


def add_chips(slide, chip_data):
    """Three chips in the lower-middle area for quick exec scan."""
    x = 0.3
    for label, value in chip_data[:3]:
        chip = slide.shapes.add_shape(1, Inches(x), Inches(4.9), Inches(3.65), Inches(0.58))
        chip.fill.solid()
        chip.fill.fore_color.rgb = BOSCH_LIGHT
        chip.line.fill.background()

        lb = slide.shapes.add_textbox(Inches(x + 0.12), Inches(4.96), Inches(1.8), Inches(0.18))
        lp = lb.text_frame.paragraphs[0]
        lp.text = label
        lp.font.size = Pt(10)
        lp.font.bold = True
        lp.font.color.rgb = MUTED

        vb = slide.shapes.add_textbox(Inches(x + 0.12), Inches(5.14), Inches(3.35), Inches(0.24))
        vp = vb.text_frame.paragraphs[0]
        vp.text = value
        vp.font.size = Pt(13)
        vp.font.bold = True
        vp.font.color.rgb = BOSCH_NAVY
        x += 3.84


def add_content(slide, bullets):
    cp = None
    for sh in slide.placeholders:
        if sh.placeholder_format.idx == 1:
            cp = sh
            break
    if cp is None:
        cp = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.8))

    tf = cp.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(22 if i == 0 else 18)
        p.font.bold = (i == 0)
        p.font.color.rgb = TEXT_DARK


# ─────────────────────────────────────────────────────────────
# Slide builders
# ─────────────────────────────────────────────────────────────

def add_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_ph(slide, 0, "Project Trinity-CAM")
    subtitle = (
        "Management Stakeholder Presentation\n"
        "JCI Aircondition Carve-Out into Robert Bosch GmbH\n"
        f"Integration Model  |  Report Date: {REPORT_DATE}"
    )
    if not set_ph(slide, 1, subtitle):
        box = slide.shapes.add_textbox(Inches(1.0), Inches(4.9), Inches(10.5), Inches(1.2))
        box.text_frame.text = subtitle


def add_content_slide(prs: Presentation, title: str, subtitle: str,
                      bullets: list, chips: list = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[10])

    if not set_ph(slide, 0, title):
        t = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(10.0), Inches(0.5))
        t.text_frame.text = title
    if not set_ph(slide, 15, subtitle):
        st = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.0), Inches(0.5))
        st.text_frame.text = subtitle

    add_content(slide, bullets)
    add_accent_bar(slide)
    add_corner_dot(slide)
    if chips:
        add_chips(slide, chips)


def add_timeline_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[10])

    set_ph(slide, 0, "Programme Timeline & Quality Gates")
    if not set_ph(slide, 15, "Jul 2026 — Apr 2028  |  6 gates  |  Integration model milestone plan"):
        st = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.0), Inches(0.5))
        st.text_frame.text = "Jul 2026 — Apr 2028  |  6 gates  |  Integration model milestone plan"

    # Timeline base line
    line = slide.shapes.add_shape(1, Inches(0.6), Inches(2.85), Inches(11.8), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = BOSCH_LIGHT
    line.line.fill.background()

    milestones = [
        ("QG0",    "1 Jul 26",  0.6,  BOSCH_BLUE),
        ("QG1",    "1 Oct 26",  2.9,  BOSCH_MID),
        ("QG2&3",  "31 Jul 27", 5.8,  BOSCH_MID),
        ("QG4",    "8 Dec 27",  9.0,  BOSCH_BLUE),
        ("GoLive", "1 Jan 28",  10.0, GREEN),
        ("QG5",    "1 Apr 28",  11.8, RGBColor(0x35, 0x7A, 0xB7)),
    ]

    for label, dt, x_pos, color in milestones:
        # circle
        circle = slide.shapes.add_shape(9, Inches(x_pos - 0.18), Inches(2.65),
                                        Inches(0.36), Inches(0.36))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.color.rgb = WHITE
        circle.line.width = Pt(2)

        # label
        lb = slide.shapes.add_textbox(Inches(x_pos - 0.42), Inches(3.1), Inches(0.84), Inches(0.28))
        lp = lb.text_frame.paragraphs[0]
        lp.text = label
        lp.alignment = PP_ALIGN.CENTER
        lp.font.size = Pt(12)
        lp.font.bold = True
        lp.font.color.rgb = BOSCH_NAVY

        # date
        db = slide.shapes.add_textbox(Inches(x_pos - 0.48), Inches(3.38), Inches(0.96), Inches(0.25))
        dp = db.text_frame.paragraphs[0]
        dp.text = dt
        dp.alignment = PP_ALIGN.CENTER
        dp.font.size = Pt(9)
        dp.font.color.rgb = MUTED

    # Phase annotations above line
    phases = [
        (0.6,  2.2, "Ph1 Initiation", BOSCH_NAVY),
        (2.9,  2.2, "Ph2 MZ Build", BOSCH_MID),
        (5.8,  2.2, "Ph3 Test & Migrate", BOSCH_MID),
        (9.0,  2.2, "Ph4 Final Ready", BOSCH_NAVY),
        (10.4, 2.2, "Ph5 Hypercare", MUTED),
    ]
    for px, py, txt, col in phases:
        pb = slide.shapes.add_textbox(Inches(px - 0.1), Inches(py),
                                      Inches(2.0), Inches(0.25))
        pp = pb.text_frame.paragraphs[0]
        pp.text = txt
        pp.font.size = Pt(9)
        pp.font.italic = True
        pp.font.color.rgb = col

    # Notes box
    notes = slide.shapes.add_textbox(Inches(0.6), Inches(3.8), Inches(11.8), Inches(1.6))
    tf = notes.text_frame
    tf.clear()
    note_lines = [
        "QG0 Jul 2026: PMO + Infosys mobilised. Initiation phase kickoff.",
        "QG1 Oct 2026: Concept approved. App inventory, MZ arch, CAPEX budget signed off.",
        "QG2&3 Jul 2027: Merger Zone DC live. SAP copy done. Wave 1 (~400 apps) validated.",
        "QG4 Dec 2027: All 12,000 users + 1,800+ apps on MZ. SAP Mock 2 passed. UAT signed.",
        "GoLive Jan 2028: Day 1 cutover. MZ live. TSA dependencies fully exited. Hypercare active.",
        "QG5 Apr 2028: 90-day hypercare complete. Bosch IT handover. Programme closed.",
    ]
    for i, txt in enumerate(note_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(4)

    add_accent_bar(slide)
    add_corner_dot(slide)
    add_chips(slide, [("GoLive", "1 Jan 2028"), ("Duration", "21 months"), ("Gates", "QG0-QG5")])


# ─────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────

def main():
    prs = create_prs()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    clear_slides(prs)

    # --- COVER ---
    add_cover(prs)

    # --- SLIDE 2: Executive Summary ---
    add_content_slide(
        prs,
        title    = "Executive Summary",
        subtitle = "Trinity-CAM — Programme Status Overview  |  Integration Model",
        bullets  = [
            "Programme formally commenced 1 Jul 2026. KPMG PMO mobilised. Infosys SOW signed.",
            "Scope: 48 sites, 12,000 users, 1,800+ applications including major SAP landscape.",
            "Integration Model: JCI IT assets migrate to Infosys-operated Merger Zone, then Bosch IT.",
            "Overall RAG: AMBER. SPI 1.00 / CPI 1.00 at QG0 baseline. No schedule or cost overrun.",
            "Key risk R001 (SAP landscape complexity, score 20) requires response plan by Aug 2026 SteerCo.",
            "GoLive target confirmed: 1 January 2028. JCI TSA expired 30 Jun 2026 — no extension.",
        ],
        chips = [("GoLive", "1 Jan 2028"), ("Overall RAG", "AMBER"), ("Top Risk", "Score 20 (SAP)")],
    )

    # --- SLIDE 3: Scope and Integration Model ---
    add_content_slide(
        prs,
        title    = "Scope & Integration Model",
        subtitle = "What is in scope and the two-hop migration path",
        bullets  = [
            "Migration path: JCI IT (source) → Merger Zone, Infosys-operated (transient) → Bosch IT (target).",
            "In scope: 12,000 users, all devices, M365 tenant, 1,800+ applications, full SAP landscape.",
            "Infrastructure: New Merger Zone DC/cloud (Infosys build); SD-WAN to 48 global sites.",
            "SAP: Full system copy JCI → MZ; client separation; 2 mock cutovers before GoLive.",
            "3 application migration waves + dedicated SAP track; GDPR-compliant 48-jurisdiction data transfers.",
            "Out of scope: Bosch standard platform harmonisation; SAP feature changes; hardware refresh.",
        ],
        chips = [("Sites", "48 Global"), ("Users", "12,000"), ("Applications", "1,800+")],
    )

    # --- SLIDE 4: Timeline and Quality Gates ---
    add_timeline_slide(prs)

    # --- SLIDE 5: Workstream Status & Readiness ---
    add_content_slide(
        prs,
        title    = "Workstream Status & Readiness",
        subtitle = "9 workstreams across KPMG and Infosys delivery",
        bullets  = [
            "GREEN: Programme Control (KPMG PMO), End-User Workplace, Identity & Access, Data Migration.",
            "AMBER: MZ Infrastructure (MZ arch pending QG1), SAP Migration (R001 response in progress).",
            "AMBER: Application Migration (1,800+ inventory 85% complete), Security & GDPR (DPIAs in progress).",
            "AMBER: TSA Exit & HR/Legal (TUPE works council notifications filed in DE; FR Aug 2026).",
            "Overall readiness: 12% (Initiation phase). Target 100% at QG4 (8 Dec 2027).",
            "No RED workstreams at programme baseline. Monitoring R001 SAP and R006 GDPR.",
        ],
        chips = [("GREEN WS", "4/9"), ("AMBER WS", "5/9"), ("RED WS", "0/9")],
    )

    # --- SLIDE 6: Budget Overview ---
    add_content_slide(
        prs,
        title    = "Budget Overview",
        subtitle = "Labour cost plan | CAPEX TBC at QG1",
        bullets  = [
            "Total Labour Budget: EUR 7,873,600 (KPMG + Infosys). Contingency +15%: EUR 9,054,640.",
            "CAPEX (MZ infrastructure, M365 licences, 48-site SD-WAN): TBC at QG1 — Infosys to submit.",
            "KPMG share (PMO + Architecture + SAP Advisory): EUR 2,562,400 (32.5% of labour).",
            "Infosys share (PM, Infra, IAM, SAP, Apps, Data, Service): EUR 5,311,200 (67.5% of labour).",
            "Phase 2 MZ Build: EUR 2.76M (35% of total) — largest programme phase by cost.",
            "CPI 1.00 at QG0 baseline. No variance to date. Budget release gate: QG1 approval.",
        ],
        chips = [("Labour Total", "EUR 7.87M"), ("Contingency", "EUR 1.18M"), ("CAPEX", "TBC QG1")],
    )

    # --- SLIDE 7: Top Risks and Mitigations ---
    add_content_slide(
        prs,
        title    = "Top Risks & Mitigations",
        subtitle = "From Trinity-CAM Risk Register — 25 risks (24 threats, 1 opportunity)",
        bullets  = [
            "R001 SAP Complexity (score 20 — HIGH): System copy delay risk. Response plan tabled Aug 2026 SteerCo.",
            "R003 SAP Mock Cutover 2 timing (score 15 — HIGH): Mock 1 (Jul 2027) validates remediation window.",
            "R002 Infosys MZ delivery delay (score 12): Cloud fallback option being assessed; procurement at QG1.",
            "R006 GDPR breach exposure (score 10): DPIA active across 48 jurisdictions; legal counsel engaged.",
            "R005 TUPE non-compliance DE/FR (score 10): Works council Germany notified; France EWC Aug 2026.",
            "R025 OPPTY: MZ cloud modernisation (score 12): EUR 3-5M/yr Bosch OPEX saving; Infosys cloud lead.",
        ],
        chips = [("High risks", "2 (R001, R003)"), ("Med risks", "22 open"), ("Opportunity", "1 (R025)")],
    )

    # --- SLIDE 8: Management Decisions Required ---
    add_content_slide(
        prs,
        title    = "Management Decisions Required",
        subtitle = "Actions required from Bosch Steering Committee and JCI",
        bullets  = [
            "DECISION 1: Approve QG1 CAPEX envelope (MZ infrastructure, M365, SD-WAN) — target vote Oct 2026.",
            "DECISION 2: Formally adopt Infosys cloud fallback option into MZ architecture (risk R002 mitigation).",
            "DECISION 3: Confirm Bosch-side IT integration scope and resources to be ready post-GoLive (Jan 2028).",
            "DECISION 4: Mandate JCI senior IT cooperation SLA for application inventory and data quality.",
            "DECISION 5: Approve August 2026 SteerCo date to review R001 SAP risk response plan.",
            "DECISION 6: Confirm GoLive date 1 Jan 2028 as hard business deadline — no further flex.",
        ],
        chips = [("Decisions open", "6"), ("Next SteerCo", "Aug 2026"), ("GoLive confirm", "1 Jan 2028")],
    )

    prs.save(str(OUTPUT))
    print(f"[Trinity-CAM] Stakeholder Presentation: {OUTPUT}")


if __name__ == "__main__":
    main()
